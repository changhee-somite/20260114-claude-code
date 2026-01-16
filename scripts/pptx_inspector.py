#!/usr/bin/env python3
"""
PPTX Inspector - Validate PowerPoint presentations for common issues.

This tool detects layout breakage and structural issues that may not be
visible through the python-pptx object model but cause rendering problems
in PowerPoint.

Key validations:
- Level 1: Quick object property checks (dimensions, positions)
- Level 2: XML structure validation (the critical position offset bug)
- Level 3: Content validation (placeholder text, text length)
- Level 4: Layout overlap detection

Usage:
    python pptx_inspector.py presentation.pptx
    python pptx_inspector.py presentation.pptx --level 4
    python pptx_inspector.py presentation.pptx --json
    python pptx_inspector.py presentation.pptx --fail-on warning

Author: Claude Code
Date: 2026-01-16
"""

import argparse
import json
import re
import sys
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def inspect_level1(pptx_path):
    """Quick inspection using python-pptx object model."""
    issues = []
    prs = Presentation(str(pptx_path))

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1

        for shape in slide.shapes:
            shape_name = shape.name

            # Check 1: Zero dimensions
            if shape.width == 0:
                issues.append({
                    'slide': slide_num,
                    'shape': shape_name,
                    'type': 'ZERO_WIDTH',
                    'severity': 'ERROR',
                    'description': f'{shape_name} has width=0'
                })

            if shape.height == 0:
                issues.append({
                    'slide': slide_num,
                    'shape': shape_name,
                    'type': 'ZERO_HEIGHT',
                    'severity': 'ERROR',
                    'description': f'{shape_name} has height=0'
                })

            # Check 2: Negative positions
            if shape.left < 0:
                issues.append({
                    'slide': slide_num,
                    'shape': shape_name,
                    'type': 'NEGATIVE_LEFT',
                    'severity': 'WARNING',
                    'description': f'{shape_name} has negative left position'
                })

            if shape.top < 0:
                issues.append({
                    'slide': slide_num,
                    'shape': shape_name,
                    'type': 'NEGATIVE_TOP',
                    'severity': 'WARNING',
                    'description': f'{shape_name} has negative top position'
                })

            # Check 3: Shape extends beyond slide bounds (with 10% tolerance)
            if shape.left + shape.width > slide_width * 1.1:
                issues.append({
                    'slide': slide_num,
                    'shape': shape_name,
                    'type': 'EXTENDS_RIGHT',
                    'severity': 'WARNING',
                    'description': f'{shape_name} extends beyond slide right edge'
                })

            if shape.top + shape.height > slide_height * 1.1:
                issues.append({
                    'slide': slide_num,
                    'shape': shape_name,
                    'type': 'EXTENDS_BOTTOM',
                    'severity': 'WARNING',
                    'description': f'{shape_name} extends beyond slide bottom edge'
                })

    return issues


def inspect_level2_xml(pptx_path):
    """
    Deep inspection of XML structure.

    This catches the CRITICAL position offset bug where python-pptx creates
    <a:xfrm> with only <a:ext> (size) but not <a:off> (position).
    """
    issues = []

    with zipfile.ZipFile(str(pptx_path), 'r') as z:
        # Find all slide XML files
        slide_files = [f for f in z.namelist() if re.match(r'ppt/slides/slide\d+\.xml', f)]

        for slide_file in sorted(slide_files):
            slide_num = int(re.search(r'slide(\d+)', slide_file).group(1))

            with z.open(slide_file) as f:
                content = f.read().decode('utf-8')

                # Check 1: xfrm with ext but without off (THE CRITICAL BUG)
                xfrm_blocks = re.findall(r'<a:xfrm>.*?</a:xfrm>', content, re.DOTALL)

                for block in xfrm_blocks:
                    has_off = '<a:off' in block
                    has_ext = '<a:ext' in block

                    if has_ext and not has_off:
                        # Extract size for context
                        ext_match = re.search(r'<a:ext cx="(\d+)" cy="(\d+)"', block)
                        size_info = ""
                        if ext_match:
                            w = int(ext_match.group(1)) / 914400
                            h = int(ext_match.group(2)) / 914400
                            size_info = f" (size: {w:.2f}\" x {h:.2f}\")"

                        issues.append({
                            'slide': slide_num,
                            'type': 'MISSING_POSITION_OFFSET',
                            'severity': 'CRITICAL',
                            'description': f'xfrm has size but no position{size_info} - will cause rendering issues'
                        })

                # Check 2: Empty embed references
                if 'r:embed=""' in content:
                    issues.append({
                        'slide': slide_num,
                        'type': 'EMPTY_EMBED_REFERENCE',
                        'severity': 'ERROR',
                        'description': 'Empty embed reference - image may not display'
                    })

        # Check relationship files for missing media
        for slide_file in slide_files:
            slide_num = int(re.search(r'slide(\d+)', slide_file).group(1))
            rels_file = slide_file.replace('slides/', 'slides/_rels/') + '.rels'

            try:
                with z.open(rels_file) as f:
                    rels_content = f.read().decode('utf-8')
                    media_refs = re.findall(r'Target="\.\./media/([^"]+)"', rels_content)
                    media_files = [f for f in z.namelist() if f.startswith('ppt/media/')]

                    for ref in media_refs:
                        expected = f'ppt/media/{ref}'
                        if expected not in media_files:
                            issues.append({
                                'slide': slide_num,
                                'type': 'MISSING_MEDIA_FILE',
                                'severity': 'ERROR',
                                'description': f'Referenced media file not found: {ref}'
                            })
            except KeyError:
                pass  # No rels file is OK for some slides

    return issues


def inspect_level3_content(pptx_path):
    """Validate content integrity."""
    issues = []
    prs = Presentation(str(pptx_path))

    placeholder_patterns = [
        (r'\[Figure:.*?\]', 'Figure reference not replaced'),
        (r'\[Image:.*?\]', 'Image reference not replaced'),
        (r'Lorem ipsum', 'Lorem ipsum placeholder text'),
        (r'Click to add', 'Unreplaced placeholder prompt'),
        (r'Title Text', 'Unreplaced title placeholder'),
        (r'Body Level', 'Unreplaced body placeholder'),
    ]

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1

        for shape in slide.shapes:
            if not hasattr(shape, 'text_frame') or not shape.has_text_frame:
                continue

            text = shape.text_frame.text

            # Check 1: Placeholder text not replaced
            for pattern, description in placeholder_patterns:
                if re.search(pattern, text, re.IGNORECASE):
                    issues.append({
                        'slide': slide_num,
                        'shape': shape.name,
                        'type': 'UNREPLACED_PLACEHOLDER',
                        'severity': 'WARNING',
                        'description': f'{description}: "{text[:50]}..."'
                    })
                    break

            # Check 2: Extremely long text
            if len(text) > 2000:
                issues.append({
                    'slide': slide_num,
                    'shape': shape.name,
                    'type': 'VERY_LONG_TEXT',
                    'severity': 'WARNING',
                    'description': f'Text length ({len(text)} chars) may overflow'
                })

            # Check 3: Whitespace only
            if text and not text.strip():
                issues.append({
                    'slide': slide_num,
                    'shape': shape.name,
                    'type': 'WHITESPACE_ONLY',
                    'severity': 'INFO',
                    'description': 'Shape contains only whitespace'
                })

    return issues


def inspect_level4_overlaps(pptx_path, tolerance=0.1):
    """Detect overlapping shapes that might cause visual issues."""
    issues = []
    prs = Presentation(str(pptx_path))

    def get_bounds(shape):
        """Get shape bounds in inches."""
        l = shape.left / 914400
        t = shape.top / 914400
        r = (shape.left + shape.width) / 914400
        b = (shape.top + shape.height) / 914400
        return l, t, r, b

    def shapes_overlap(s1, s2):
        """Check if two shapes overlap."""
        l1, t1, r1, b1 = get_bounds(s1)
        l2, t2, r2, b2 = get_bounds(s2)

        h_overlap = not (r1 <= l2 + tolerance or l1 >= r2 - tolerance)
        v_overlap = not (b1 <= t2 + tolerance or t1 >= b2 - tolerance)

        return h_overlap and v_overlap

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        shapes = list(slide.shapes)

        for j, s1 in enumerate(shapes):
            for s2 in shapes[j+1:]:
                # Skip slide number placeholders
                if 'Slide Number' in s1.name or 'Slide Number' in s2.name:
                    continue

                if shapes_overlap(s1, s2):
                    s1_has_text = (hasattr(s1, 'text_frame') and
                                  s1.has_text_frame and
                                  s1.text_frame.text.strip())
                    s2_has_text = (hasattr(s2, 'text_frame') and
                                  s2.has_text_frame and
                                  s2.text_frame.text.strip())

                    # Text overlapping text is a warning
                    if s1_has_text and s2_has_text:
                        issues.append({
                            'slide': slide_num,
                            'type': 'TEXT_OVERLAP',
                            'severity': 'WARNING',
                            'description': f'"{s1.name}" overlaps with "{s2.name}" - both have text',
                            'shapes': [s1.name, s2.name]
                        })
                    # Image overlapping text placeholder might be intentional
                    elif (s1.shape_type == MSO_SHAPE_TYPE.PICTURE or
                          s2.shape_type == MSO_SHAPE_TYPE.PICTURE):
                        # Only report if significant overlap
                        pass  # Often intentional

    return issues


def main():
    parser = argparse.ArgumentParser(
        description='Inspect PPTX for layout issues and structural problems',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Inspection Levels:
  1 - Quick: Object property checks (dimensions, positions)
  2 - XML: Structure validation (position offset bug) [DEFAULT]
  3 - Content: Text validation (placeholders, length)
  4 - Full: All above plus overlap detection

Examples:
  python pptx_inspector.py presentation.pptx
  python pptx_inspector.py presentation.pptx --level 4
  python pptx_inspector.py presentation.pptx --json --fail-on warning
        """
    )
    parser.add_argument('pptx_file', help='PowerPoint file to inspect')
    parser.add_argument('--level', '-l', type=int, default=2, choices=[1, 2, 3, 4],
                       help='Inspection depth (default: 2)')
    parser.add_argument('--json', '-j', action='store_true',
                       help='Output as JSON')
    parser.add_argument('--fail-on', '-f',
                       choices=['critical', 'error', 'warning', 'info'],
                       default='critical',
                       help='Exit with error if issues at this level or above (default: critical)')
    parser.add_argument('--quiet', '-q', action='store_true',
                       help='Only output if issues found')

    args = parser.parse_args()

    pptx_path = Path(args.pptx_file)
    if not pptx_path.exists():
        print(f"Error: File not found: {pptx_path}", file=sys.stderr)
        sys.exit(1)

    all_issues = []

    # Run inspections up to requested level
    if args.level >= 1:
        all_issues.extend(inspect_level1(pptx_path))
    if args.level >= 2:
        all_issues.extend(inspect_level2_xml(pptx_path))
    if args.level >= 3:
        all_issues.extend(inspect_level3_content(pptx_path))
    if args.level >= 4:
        all_issues.extend(inspect_level4_overlaps(pptx_path))

    # Output
    if args.json:
        print(json.dumps(all_issues, indent=2))
    else:
        if not all_issues:
            if not args.quiet:
                print(f"✅ {pptx_path.name}: No issues found (level {args.level} inspection)")
        else:
            severity_order = {'CRITICAL': 0, 'ERROR': 1, 'WARNING': 2, 'INFO': 3}
            sorted_issues = sorted(
                all_issues,
                key=lambda x: (severity_order.get(x.get('severity', 'INFO'), 4),
                              x.get('slide', 0))
            )

            print(f"{'='*60}")
            print(f"PPTX Inspector: {pptx_path.name}")
            print(f"{'='*60}")
            print(f"Found {len(all_issues)} issue(s):\n")

            for issue in sorted_issues:
                sev = issue.get('severity', 'INFO')
                icon = {
                    'CRITICAL': '🔴',
                    'ERROR': '❌',
                    'WARNING': '⚠️',
                    'INFO': 'ℹ️'
                }.get(sev, '?')

                slide = issue.get('slide', '?')
                desc = issue.get('description', str(issue))
                print(f"{icon} [{sev}] Slide {slide}: {desc}")

            print()

    # Determine exit code
    fail_levels = {
        'critical': ['CRITICAL'],
        'error': ['CRITICAL', 'ERROR'],
        'warning': ['CRITICAL', 'ERROR', 'WARNING'],
        'info': ['CRITICAL', 'ERROR', 'WARNING', 'INFO']
    }

    should_fail = any(
        issue.get('severity') in fail_levels[args.fail_on]
        for issue in all_issues
        if isinstance(issue, dict)
    )

    sys.exit(1 if should_fail else 0)


if __name__ == '__main__':
    main()
