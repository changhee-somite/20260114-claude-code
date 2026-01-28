#!/usr/bin/env python3
"""
Table Insertion Script

This script adds tables to PowerPoint presentations based on a JSON mapping file.
It follows the same pattern as add_images_only.py and handles the python-pptx
position offset bug correctly.

Usage:
    python add_tables.py input.pptx output.pptx --mapping table-mapping.json

Author: Claude Code
Date: 2026-01-18
"""

import argparse
import json
import re
import sys
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import PP_PLACEHOLDER

from table_renderer import (
    render_table, parse_markdown_table, TableConfig, TableStyle
)


def find_body_placeholder(slide):
    """Find the body/content placeholder."""
    for shape in slide.placeholders:
        ph_type = shape.placeholder_format.type
        if ph_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
            return shape
    return None


def resize_placeholder_width(shape, new_width_inches):
    """
    Resize placeholder width while preserving height AND position.

    CRITICAL: When setting width on a placeholder, python-pptx creates an
    <a:xfrm> element with only <a:ext> (size) but not <a:off> (position).
    This breaks PowerPoint's layout inheritance and causes visual glitches.

    Solution: Explicitly set left, top, width, height together.
    """
    original_left = shape.left
    original_top = shape.top
    original_height = shape.height

    shape.left = original_left
    shape.top = original_top
    shape.width = Inches(new_width_inches)
    shape.height = original_height

    return shape


def validate_presentation(pptx_path, table_slides):
    """
    Validate the generated presentation for common issues.
    Returns list of issues found.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    issues = []
    prs = Presentation(str(pptx_path))

    for slide_num in table_slides:
        slide_idx = slide_num - 1
        if slide_idx >= len(prs.slides):
            continue

        slide = prs.slides[slide_idx]

        # Check 1: Table presence
        has_table = any(s.shape_type == MSO_SHAPE_TYPE.TABLE for s in slide.shapes)
        if not has_table:
            issues.append(f"Slide {slide_num}: No table found")

        # Check 2: Table dimensions
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table = shape.table
                if len(table.rows) == 0:
                    issues.append(f"Slide {slide_num}: Table has no rows")
                if len(table.columns) == 0:
                    issues.append(f"Slide {slide_num}: Table has no columns")

                # Check for empty cells in header
                if len(table.rows) > 0:
                    header_row = table.rows[0]
                    empty_headers = sum(1 for cell in header_row.cells if not cell.text.strip())
                    if empty_headers > 0:
                        issues.append(f"Slide {slide_num}: Table has {empty_headers} empty header cell(s)")

        # Check 3: Text placeholder dimensions
        for shape in slide.shapes:
            if 'Text Placeholder' in shape.name:
                height = shape.height / 914400
                width = shape.width / 914400

                if height < 0.5:
                    issues.append(f"Slide {slide_num}: Text placeholder height too small ({height:.2f}\")")
                if width < 2.0:
                    issues.append(f"Slide {slide_num}: Text placeholder width too small ({width:.2f}\")")

    # Check 4: XML structure (position offsets)
    with zipfile.ZipFile(str(pptx_path), 'r') as z:
        for slide_num in table_slides:
            slide_xml = f'ppt/slides/slide{slide_num}.xml'
            try:
                with z.open(slide_xml) as f:
                    content = f.read().decode('utf-8')

                    xfrm_blocks = re.findall(r'<a:xfrm>.*?</a:xfrm>', content, re.DOTALL)
                    for block in xfrm_blocks:
                        if '<a:ext' in block and '<a:off' not in block:
                            issues.append(
                                f"Slide {slide_num}: xfrm missing position offset (will cause display issues)"
                            )
                            break
            except KeyError:
                pass

    return issues


def main():
    parser = argparse.ArgumentParser(description="Add tables to presentation")
    parser.add_argument("input_pptx", help="Input PowerPoint file")
    parser.add_argument("output_pptx", help="Output PowerPoint file")
    parser.add_argument("--mapping", "-m", required=True, help="Table mapping JSON file")
    parser.add_argument("--text-width", type=float, default=6.5,
                        help="Width for text placeholder on table slides")
    parser.add_argument("--table-left", type=float, default=7.0,
                        help="Left position for tables")
    parser.add_argument("--table-top", type=float, default=1.5,
                        help="Top position for tables")
    parser.add_argument("--table-width", type=float, default=5.5,
                        help="Table width")

    args = parser.parse_args()

    input_path = Path(args.input_pptx)
    output_path = Path(args.output_pptx)
    mapping_path = Path(args.mapping)

    if not input_path.exists():
        print(f"Error: Input file not found: {input_path}")
        sys.exit(1)

    if not mapping_path.exists():
        print(f"Error: Mapping file not found: {mapping_path}")
        sys.exit(1)

    # Load mapping
    with open(mapping_path, 'r') as f:
        mapping = json.load(f)

    # Load presentation
    prs = Presentation(str(input_path))

    print(f"Input: {input_path}")
    print(f"Slides: {len(prs.slides)}")
    print(f"Mapping: {mapping_path}")
    print()

    # Process each slide in mapping
    table_count = 0

    for slide_config in mapping.get('slides', []):
        slide_num = slide_config.get('slide_number')
        if slide_num is None:
            continue

        slide_idx = slide_num - 1
        if slide_idx < 0 or slide_idx >= len(prs.slides):
            print(f"Warning: Slide {slide_num} out of range")
            continue

        slide = prs.slides[slide_idx]

        print(f"Slide {slide_num}:")

        # Step 1: Resize text placeholder (if exists)
        body = find_body_placeholder(slide)
        if body:
            old_width = body.width / 914400
            resize_placeholder_width(body, args.text_width)
            new_width = body.width / 914400
            print(f"  Text placeholder: {old_width:.2f}\" → {new_width:.2f}\"")

        # Step 2: Insert tables
        for table_config in slide_config.get('tables', []):
            # Get table data
            headers = table_config.get('headers', [])
            rows = table_config.get('rows', [])
            markdown = table_config.get('markdown')

            # Parse markdown if provided
            if markdown:
                parsed = parse_markdown_table(markdown)
                headers = parsed.get('headers', headers)
                rows = parsed.get('rows', rows)

            if not headers:
                print("  Warning: No headers for table, skipping")
                continue

            # Position configuration
            position = table_config.get('position', {})
            config = TableConfig(
                left=position.get('left', args.table_left),
                top=position.get('top', args.table_top),
                width=position.get('width', args.table_width),
                row_height=position.get('row_height', 0.4)
            )

            # Style configuration
            style_config = table_config.get('style', {})
            style = TableStyle(
                header_fill=style_config.get('header_fill', '#4472C4'),
                header_text=style_config.get('header_text', '#FFFFFF'),
                row_fill=style_config.get('row_fill', '#FFFFFF'),
                alt_row_fill=style_config.get('alt_row_fill', '#F2F2F2'),
                text_color=style_config.get('text_color', '#000000'),
                font_size=style_config.get('font_size', 10),
                header_font_size=style_config.get('header_font_size', 11),
                use_alt_rows=style_config.get('use_alt_rows', True)
            )

            # Render table
            spec = {"headers": headers, "rows": rows}
            table_shape = render_table(slide, spec, config, style)

            if table_shape:
                table_count += 1
                tbl_left = table_shape.left / 914400
                tbl_top = table_shape.top / 914400
                tbl_width = table_shape.width / 914400
                tbl_height = table_shape.height / 914400
                print(f"  Table: {len(headers)} cols x {len(rows)+1} rows")
                print(f"    Position: ({tbl_left:.2f}\", {tbl_top:.2f}\")")
                print(f"    Size: {tbl_width:.2f}\" x {tbl_height:.2f}\"")

        print()

    # Save
    prs.save(str(output_path))

    print(f"Inserted {table_count} table(s)")
    print(f"Saved to: {output_path}")

    # Validation
    print()
    print("VALIDATION")
    print("=" * 50)

    table_slides = [s.get('slide_number') for s in mapping.get('slides', []) if s.get('slide_number')]
    issues = validate_presentation(output_path, table_slides)

    if issues:
        print(f"❌ Found {len(issues)} issue(s):")
        for issue in issues:
            print(f"  - {issue}")
        sys.exit(1)
    else:
        print("✅ All validation checks passed!")


if __name__ == "__main__":
    main()
