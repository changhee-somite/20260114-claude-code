#!/usr/bin/env python3
"""
Minimal Image Insertion Script

This script ONLY:
1. Resizes text placeholders on slides that will have images
2. Inserts images

It does NOT modify any text content - that should already be correct
from the replace.py step.

Usage:
    python add_images_only.py input.pptx output.pptx --mapping image-mapping.json

Author: Claude Code
Date: 2026-01-16
"""

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import PP_PLACEHOLDER
from PIL import Image


def find_body_placeholder(slide):
    """Find the body/content placeholder."""
    for shape in slide.placeholders:
        ph_type = shape.placeholder_format.type
        if ph_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
            return shape
    return None


def resize_placeholder_width(shape, new_width_inches):
    """
    Resize placeholder width while preserving height AND position.

    CRITICAL: When setting width on a placeholder, python-pptx creates an
    <a:xfrm> element with only <a:ext> (size) but not <a:off> (position).
    This breaks PowerPoint's layout inheritance and causes visual glitches.

    Solution: Explicitly set left, top, width, height together.
    """
    # Save original values
    original_left = shape.left
    original_top = shape.top
    original_height = shape.height

    # Set all four properties explicitly
    # This ensures python-pptx creates complete <a:xfrm> with both <a:off> and <a:ext>
    shape.left = original_left
    shape.top = original_top
    shape.width = Inches(new_width_inches)
    shape.height = original_height

    return shape


def insert_image(slide, image_path, left, top, max_width, max_height):
    """Insert image with proper aspect ratio."""

    if not image_path.exists():
        print(f"    WARNING: Image not found: {image_path}")
        return None

    # Get image dimensions
    with Image.open(image_path) as img:
        img_w, img_h = img.size

    aspect = img_w / img_h

    # Calculate size to fit within bounds
    if aspect > max_width / max_height:
        # Width-constrained
        width = max_width
        height = width / aspect
    else:
        # Height-constrained
        height = max_height
        width = height * aspect

    pic = slide.shapes.add_picture(
        str(image_path),
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height)
    )

    return pic


def validate_presentation(pptx_path, image_slides):
    """
    Validate the generated presentation for common issues.
    Returns list of issues found.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    import zipfile
    import xml.etree.ElementTree as ET

    issues = []
    prs = Presentation(str(pptx_path))

    for slide_num in image_slides:
        slide_idx = slide_num - 1
        if slide_idx >= len(prs.slides):
            continue

        slide = prs.slides[slide_idx]

        # Check 1: Placeholder dimensions
        for shape in slide.shapes:
            if 'Text Placeholder' in shape.name:
                height = shape.height / 914400
                width = shape.width / 914400

                if height < 0.5:
                    issues.append(f"Slide {slide_num}: Text placeholder height too small ({height:.2f}\")")

                if width < 2.0:
                    issues.append(f"Slide {slide_num}: Text placeholder width too small ({width:.2f}\")")

        # Check 2: Image presence
        has_image = any(s.shape_type == MSO_SHAPE_TYPE.PICTURE for s in slide.shapes)
        if not has_image:
            issues.append(f"Slide {slide_num}: No image found")

    # Check 3: XML structure (position offsets)
    with zipfile.ZipFile(str(pptx_path), 'r') as z:
        for slide_num in image_slides:
            slide_xml = f'ppt/slides/slide{slide_num}.xml'
            try:
                with z.open(slide_xml) as f:
                    content = f.read().decode('utf-8')

                    # Check for xfrm with ext but without off
                    if '<a:xfrm><a:ext' in content and '<a:xfrm><a:off' not in content.replace('</a:xfrm>', '|||').split('|||')[0]:
                        # More precise check
                        import re
                        xfrm_blocks = re.findall(r'<a:xfrm>.*?</a:xfrm>', content, re.DOTALL)
                        for block in xfrm_blocks:
                            if '<a:ext' in block and '<a:off' not in block:
                                issues.append(f"Slide {slide_num}: xfrm missing position offset (will cause display issues)")
                                break
            except KeyError:
                pass

    return issues


def main():
    parser = argparse.ArgumentParser(description="Add images to presentation (text unchanged)")
    parser.add_argument("input_pptx", help="Input PowerPoint file (with text already set)")
    parser.add_argument("output_pptx", help="Output PowerPoint file")
    parser.add_argument("--mapping", "-m", required=True, help="Image mapping JSON file")
    parser.add_argument("--text-width", type=float, default=6.5, help="Width for text on image slides")
    parser.add_argument("--image-left", type=float, default=7.0, help="Left position for images")
    parser.add_argument("--image-top", type=float, default=1.5, help="Top position for images")
    parser.add_argument("--image-max-width", type=float, default=5.5, help="Max image width")
    parser.add_argument("--image-max-height", type=float, default=5.5, help="Max image height")

    args = parser.parse_args()

    input_path = Path(args.input_pptx)
    output_path = Path(args.output_pptx)
    mapping_path = Path(args.mapping)

    if not input_path.exists():
        print(f"Error: Input file not found: {input_path}")
        sys.exit(1)

    if not mapping_path.exists():
        print(f"Error: Mapping file not found: {mapping_path}")
        sys.exit(1)

    # Load mapping
    with open(mapping_path, 'r') as f:
        mapping = json.load(f)

    # Load presentation
    prs = Presentation(str(input_path))

    print(f"Input: {input_path}")
    print(f"Slides: {len(prs.slides)}")
    print(f"Mapping: {mapping_path}")
    print()

    # Process each slide in mapping
    image_count = 0

    for slide_config in mapping.get('slides', []):
        slide_num = slide_config.get('slide_number')
        if slide_num is None:
            continue

        slide_idx = slide_num - 1
        if slide_idx < 0 or slide_idx >= len(prs.slides):
            print(f"Warning: Slide {slide_num} out of range")
            continue

        slide = prs.slides[slide_idx]

        print(f"Slide {slide_num}:")

        # Step 1: Resize text placeholder (if exists)
        body = find_body_placeholder(slide)
        if body:
            old_width = body.width / 914400
            resize_placeholder_width(body, args.text_width)
            new_width = body.width / 914400
            new_height = body.height / 914400
            print(f"  Text placeholder: {old_width:.2f}\" → {new_width:.2f}\" (height={new_height:.2f}\")")

        # Step 2: Insert images
        for img_config in slide_config.get('images', []):
            img_path_str = img_config.get('path')
            if not img_path_str:
                continue

            # Resolve path relative to mapping file
            if img_path_str.startswith('../'):
                img_path = mapping_path.parent.parent / img_path_str[3:]
            else:
                img_path = mapping_path.parent / img_path_str

            if not img_path.exists():
                # Try figures directory
                img_path = mapping_path.parent.parent / 'figures' / Path(img_path_str).name

            pic = insert_image(
                slide,
                img_path,
                args.image_left,
                args.image_top,
                args.image_max_width,
                args.image_max_height
            )

            if pic:
                image_count += 1
                pic_left = pic.left / 914400
                pic_top = pic.top / 914400
                pic_width = pic.width / 914400
                pic_height = pic.height / 914400
                print(f"  Image: {img_path.name}")
                print(f"    Position: ({pic_left:.2f}\", {pic_top:.2f}\") Size: {pic_width:.2f}\" x {pic_height:.2f}\"")

        print()

    # Save
    prs.save(str(output_path))

    print(f"Inserted {image_count} images")
    print(f"Saved to: {output_path}")

    # Validation
    print()
    print("VALIDATION")
    print("=" * 50)

    image_slides = [s.get('slide_number') for s in mapping.get('slides', []) if s.get('slide_number')]
    issues = validate_presentation(output_path, image_slides)

    if issues:
        print(f"❌ Found {len(issues)} issue(s):")
        for issue in issues:
            print(f"  - {issue}")
        sys.exit(1)
    else:
        print("✅ All validation checks passed!")


if __name__ == "__main__":
    main()
