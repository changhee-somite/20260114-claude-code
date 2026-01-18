#!/usr/bin/env python3
"""
Diagram Renderer: Create vector graphics diagrams in PowerPoint.

This module renders structured diagram specifications as PowerPoint shapes
using python-pptx. It supports cycle, flow, comparison, and hierarchy diagrams.

Usage:
    from diagram_renderer import render_diagram

    render_diagram(slide, diagram_spec, position_config)

Author: Claude Code
Date: 2026-01-17
"""

import math
from dataclasses import dataclass
from typing import List, Dict, Optional, Tuple

from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


@dataclass
class DiagramConfig:
    """Configuration for diagram positioning."""
    left: float = 7.0      # Left position in inches
    top: float = 1.5       # Top position in inches
    width: float = 5.0     # Available width in inches
    height: float = 5.0    # Available height in inches


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color string to RGBColor."""
    hex_color = hex_color.lstrip('#')
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


def get_shape_type(shape_name: str) -> MSO_SHAPE:
    """Map shape name to MSO_SHAPE enum."""
    shape_map = {
        "rectangle": MSO_SHAPE.RECTANGLE,
        "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
        "oval": MSO_SHAPE.OVAL,
        "diamond": MSO_SHAPE.DIAMOND,
        "parallelogram": MSO_SHAPE.PARALLELOGRAM,
        "chevron": MSO_SHAPE.CHEVRON,
        "pentagon": MSO_SHAPE.PENTAGON,
        "hexagon": MSO_SHAPE.HEXAGON,
        "flowchart_process": MSO_SHAPE.FLOWCHART_PROCESS,
        "flowchart_decision": MSO_SHAPE.FLOWCHART_DECISION,
        "flowchart_terminator": MSO_SHAPE.FLOWCHART_TERMINATOR,
    }
    return shape_map.get(shape_name.lower(), MSO_SHAPE.ROUNDED_RECTANGLE)


def create_styled_shape(slide, shape_type, left, top, width, height, text, style):
    """Create a shape with text and styling."""
    mso_shape = get_shape_type(shape_type)

    shape = slide.shapes.add_shape(
        mso_shape,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height)
    )

    # Set text
    shape.text = text

    # Style the text
    for paragraph in shape.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(11)
            run.font.bold = True
            if "text_color" in style:
                run.font.color.rgb = hex_to_rgb(style["text_color"])

    # Center text vertically
    shape.text_frame.word_wrap = True
    shape.text_frame.auto_size = None

    # Apply fill
    if "fill" in style:
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(style["fill"])

    # Apply line
    if "line" in style:
        shape.line.color.rgb = hex_to_rgb(style["line"])
        shape.line.width = Pt(2)

    return shape


def add_arrow_connector(slide, start_shape, end_shape, style):
    """Add an arrow connector between two shapes."""
    # Get shape centers and edges
    start_cx = start_shape.left + start_shape.width / 2
    start_cy = start_shape.top + start_shape.height / 2
    end_cx = end_shape.left + end_shape.width / 2
    end_cy = end_shape.top + end_shape.height / 2

    # Calculate direction
    dx = end_cx - start_cx
    dy = end_cy - start_cy

    # Determine connection points (edges of shapes)
    if abs(dx) > abs(dy):
        # Horizontal connection
        if dx > 0:
            start_x = start_shape.left + start_shape.width
            start_y = start_cy
            end_x = end_shape.left
            end_y = end_cy
        else:
            start_x = start_shape.left
            start_y = start_cy
            end_x = end_shape.left + end_shape.width
            end_y = end_cy
    else:
        # Vertical connection
        if dy > 0:
            start_x = start_cx
            start_y = start_shape.top + start_shape.height
            end_x = end_cx
            end_y = end_shape.top
        else:
            start_x = start_cx
            start_y = start_shape.top
            end_x = end_cx
            end_y = end_shape.top + end_shape.height

    # Add connector
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        start_x, start_y,
        end_x, end_y
    )

    # Style the connector
    if "line" in style:
        connector.line.color.rgb = hex_to_rgb(style["line"])
    connector.line.width = Pt(2)

    return connector


def render_cycle_diagram(slide, spec: dict, config: DiagramConfig) -> List:
    """
    Render a cycle diagram with shapes arranged in a loop.

    For 4 nodes, arranges them in a square pattern:
        [1] → [2]
         ↑     ↓
        [4] ← [3]
    """
    nodes = spec.get("nodes", [])
    edges = spec.get("edges", [])
    style = spec.get("style", {})

    if not nodes:
        return []

    shapes_created = []
    node_shapes = {}

    n = len(nodes)

    # Calculate layout
    box_width = 1.6
    box_height = 0.7
    spacing = 0.4

    if n == 4:
        # Square arrangement
        positions = [
            (config.left, config.top),                                          # Top-left
            (config.left + box_width + spacing * 2, config.top),                # Top-right
            (config.left + box_width + spacing * 2, config.top + box_height + spacing * 2),  # Bottom-right
            (config.left, config.top + box_height + spacing * 2),               # Bottom-left
        ]
    elif n == 3:
        # Triangle arrangement
        cx = config.left + config.width / 2 - box_width / 2
        positions = [
            (cx, config.top),                                    # Top
            (config.left + config.width - box_width, config.top + box_height + spacing),  # Right
            (config.left, config.top + box_height + spacing),    # Left
        ]
    else:
        # Arrange in a circle
        center_x = config.left + config.width / 2
        center_y = config.top + config.height / 2
        radius = min(config.width, config.height) / 2 - box_width / 2

        positions = []
        for i in range(n):
            angle = -math.pi / 2 + (2 * math.pi * i / n)  # Start from top
            x = center_x + radius * math.cos(angle) - box_width / 2
            y = center_y + radius * math.sin(angle) - box_height / 2
            positions.append((x, y))

    # Create shapes
    shape_type = style.get("shape", "rounded_rectangle")

    for i, node in enumerate(nodes):
        if i >= len(positions):
            break

        left, top = positions[i]

        shape = create_styled_shape(
            slide,
            shape_type,
            left, top,
            box_width, box_height,
            node.get("text", ""),
            style
        )

        node_shapes[node["id"]] = shape
        shapes_created.append(shape)

    # Create connectors
    for edge in edges:
        from_id = edge.get("from")
        to_id = edge.get("to")

        if from_id in node_shapes and to_id in node_shapes:
            connector = add_arrow_connector(
                slide,
                node_shapes[from_id],
                node_shapes[to_id],
                style
            )
            shapes_created.append(connector)

    return shapes_created


def render_flow_diagram(slide, spec: dict, config: DiagramConfig, vertical: bool = False) -> List:
    """Render a linear flow diagram."""
    nodes = spec.get("nodes", [])
    edges = spec.get("edges", [])
    style = spec.get("style", {})

    if not nodes:
        return []

    shapes_created = []
    node_shapes = {}

    n = len(nodes)
    box_width = 1.5
    box_height = 0.6
    spacing = 0.3

    shape_type = style.get("shape", "rounded_rectangle")

    for i, node in enumerate(nodes):
        if vertical:
            left = config.left + (config.width - box_width) / 2
            top = config.top + i * (box_height + spacing)
        else:
            left = config.left + i * (box_width + spacing)
            top = config.top + (config.height - box_height) / 2

        shape = create_styled_shape(
            slide,
            shape_type,
            left, top,
            box_width, box_height,
            node.get("text", ""),
            style
        )

        node_shapes[node["id"]] = shape
        shapes_created.append(shape)

    # Create connectors
    for edge in edges:
        from_id = edge.get("from")
        to_id = edge.get("to")

        if from_id in node_shapes and to_id in node_shapes:
            connector = add_arrow_connector(
                slide,
                node_shapes[from_id],
                node_shapes[to_id],
                style
            )
            shapes_created.append(connector)

    return shapes_created


def render_comparison_diagram(slide, spec: dict, config: DiagramConfig) -> List:
    """Render a two-column comparison diagram."""
    nodes = spec.get("nodes", [])
    style = spec.get("style", {})

    if not nodes:
        return []

    shapes_created = []

    box_width = config.width / 2 - 0.2
    box_height = 0.6
    spacing = 0.3

    # Split nodes into two columns
    left_nodes = nodes[:len(nodes)//2]
    right_nodes = nodes[len(nodes)//2:]

    shape_type = style.get("shape", "rounded_rectangle")

    # Left column
    for i, node in enumerate(left_nodes):
        shape = create_styled_shape(
            slide,
            shape_type,
            config.left,
            config.top + i * (box_height + spacing),
            box_width, box_height,
            node.get("text", ""),
            style
        )
        shapes_created.append(shape)

    # Right column
    for i, node in enumerate(right_nodes):
        shape = create_styled_shape(
            slide,
            shape_type,
            config.left + box_width + 0.4,
            config.top + i * (box_height + spacing),
            box_width, box_height,
            node.get("text", ""),
            style
        )
        shapes_created.append(shape)

    return shapes_created


def render_diagram(slide, spec: dict, config: DiagramConfig = None) -> List:
    """
    Render a diagram on a slide based on the specification.

    Args:
        slide: PowerPoint slide object
        spec: Diagram specification dict with type, nodes, edges, style
        config: DiagramConfig for positioning

    Returns:
        List of created shapes
    """
    if config is None:
        config = DiagramConfig()

    diagram_type = spec.get("type", "flow")

    # Override config with spec values if provided
    if "position" in spec:
        pos = spec["position"]
        if pos == "right":
            config.left = 7.0
        elif pos == "left":
            config.left = 0.5
        elif pos == "center":
            config.left = 3.5

    if "width" in spec:
        config.width = float(spec["width"])

    if "top" in spec:
        config.top = float(spec["top"])

    # Dispatch to appropriate renderer
    if diagram_type == "cycle":
        return render_cycle_diagram(slide, spec, config)
    elif diagram_type == "flow":
        return render_flow_diagram(slide, spec, config, vertical=False)
    elif diagram_type == "flow_vertical":
        return render_flow_diagram(slide, spec, config, vertical=True)
    elif diagram_type == "comparison":
        return render_comparison_diagram(slide, spec, config)
    else:
        print(f"  Warning: Unknown diagram type '{diagram_type}', using flow")
        return render_flow_diagram(slide, spec, config)


def parse_diagram_block(block: str) -> dict:
    """
    Parse a ```diagram code block into a specification dict.

    Expects YAML-like format:
        type: cycle
        nodes:
          - id: n1
            text: Node 1
        edges:
          - from: n1
            to: n2
        style:
          fill: "#E8F4F8"
    """
    import re

    spec = {
        "type": "flow",
        "nodes": [],
        "edges": [],
        "style": {}
    }

    lines = block.strip().split('\n')
    current_section = None
    current_item = None

    for line in lines:
        stripped = line.strip()

        if not stripped or stripped.startswith('#'):
            continue

        # Top-level keys
        if stripped.startswith('type:'):
            spec["type"] = stripped.split(':', 1)[1].strip()
        elif stripped.startswith('position:'):
            spec["position"] = stripped.split(':', 1)[1].strip()
        elif stripped.startswith('width:'):
            spec["width"] = stripped.split(':', 1)[1].strip()
        elif stripped.startswith('top:'):
            spec["top"] = stripped.split(':', 1)[1].strip()
        elif stripped == 'nodes:':
            current_section = 'nodes'
        elif stripped == 'edges:':
            current_section = 'edges'
        elif stripped == 'style:':
            current_section = 'style'
        elif stripped.startswith('- id:'):
            # New node
            current_item = {"id": stripped.split(':', 1)[1].strip()}
            if current_section == 'nodes':
                spec["nodes"].append(current_item)
        elif stripped.startswith('- from:'):
            # New edge
            current_item = {"from": stripped.split(':', 1)[1].strip()}
            if current_section == 'edges':
                spec["edges"].append(current_item)
        elif stripped.startswith('text:') and current_item:
            current_item["text"] = stripped.split(':', 1)[1].strip()
        elif stripped.startswith('to:') and current_item:
            current_item["to"] = stripped.split(':', 1)[1].strip()
        elif current_section == 'style' and ':' in stripped:
            key, value = stripped.split(':', 1)
            spec["style"][key.strip()] = value.strip().strip('"\'')

    return spec
