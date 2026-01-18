#!/usr/bin/env python3
"""
Unified Presentation Generator

Generates PowerPoint presentations from PRESENTATION.md with integrated
text, image, and diagram handling.

Features:
- Parses PRESENTATION.md to extract slide content
- Detects slides with figures, diagrams, or text-only
- Assigns appropriate layouts (split for images/diagrams, full-width for text)
- Inserts images with proper positioning
- Renders vector diagrams using python-pptx shapes
- Resizes text placeholders to avoid overlap

Usage:
    python generate_presentation.py \\
        --template template.pptx \\
        --source docs/PRESENTATION.md \\
        --figures figures/ \\
        --output final.pptx

Author: Claude Code
Date: 2026-01-17
"""

import argparse
import json
import os
import re
import sys
from dataclasses import dataclass, field
from pathlib import Path
from typing import List, Optional, Dict, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.dml.color import RGBColor
from PIL import Image

# Import diagram renderer
from diagram_renderer import render_diagram, parse_diagram_block, DiagramConfig


@dataclass
class SlideContent:
    """Represents parsed content for a single slide."""
    number: int
    title: str = ""
    content: List[str] = field(default_factory=list)
    figure_path: Optional[str] = None
    diagram_spec: Optional[Dict] = None  # Parsed diagram specification
    notes: Optional[str] = None
    slide_type: str = "text_only"  # text_only, text_image, text_diagram, section, title


@dataclass
class LayoutConfig:
    """Configuration for slide layouts."""
    text_only_idx: int = 1
    text_image_idx: int = 1  # Same as text_only if no picture placeholder
    section_idx: int = 0
    title_idx: int = 0

    # Text placeholder dimensions for image slides (in inches)
    text_width: float = 6.5  # Width of text area when image present
    image_left: float = 7.0   # Left position of image
    image_width: float = 5.5  # Width of image
    image_top: float = 1.5    # Top position of image (below title)


def parse_presentation_md(md_path: Path) -> List[SlideContent]:
    """Parse PRESENTATION.md and extract slide content."""

    with open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()

    slides = []

    # Split by slide headers (### Slide N: or ### Slide AN:)
    slide_pattern = r'###\s+Slide\s+(\d+|A\d+):\s*(.+?)(?=###\s+Slide|\Z)'
    matches = re.findall(slide_pattern, content, re.DOTALL)

    for slide_num_str, slide_content in matches:
        slide = SlideContent(
            number=int(slide_num_str) if slide_num_str.isdigit() else 100 + int(slide_num_str[1:])
        )

        # Extract title
        title_match = re.search(r'\*\*Title\*\*:\s*(.+?)(?:\n|$)', slide_content)
        if title_match:
            slide.title = title_match.group(1).strip()

        # Extract type if specified
        type_match = re.search(r'\*\*Type\*\*:\s*(.+?)(?:\n|$)', slide_content)
        if type_match:
            type_str = type_match.group(1).strip().lower()
            if 'section' in type_str:
                slide.slide_type = 'section'

        # Extract content (bullets from **Content**: section)
        content_match = re.search(
            r'\*\*Content\*\*:\s*\n((?:(?!\*\*(?:Figure|Notes|Source)\*\*).)+)',
            slide_content,
            re.DOTALL
        )
        if content_match:
            content_text = content_match.group(1)
            # Parse bullets and other content
            bullets = []
            for line in content_text.split('\n'):
                line = line.strip()
                if line.startswith(('-', '*', '•')) and not line.startswith('**'):
                    bullet_text = re.sub(r'^[-*•]\s*', '', line)
                    if bullet_text:
                        bullets.append(bullet_text)
                elif line.startswith('|'):
                    # Table row - keep as is
                    bullets.append(line)
                elif line.startswith('>'):
                    # Quote
                    quote_text = line.lstrip('> ')
                    bullets.append(f'"{quote_text}"')
                elif line.startswith('```'):
                    continue  # Skip code fence markers
                elif line and not line.startswith('**'):
                    bullets.append(line)
            slide.content = bullets

        # Extract figure path
        figure_match = re.search(
            r'\*\*Figure\*\*:\s*\[.+?\]\(([^)]+)\)',
            slide_content
        )
        if figure_match:
            slide.figure_path = figure_match.group(1)
            slide.slide_type = 'text_image'

        # Extract diagram specification
        diagram_match = re.search(
            r'\*\*Diagram\*\*:\s*\n```diagram\n(.+?)```',
            slide_content,
            re.DOTALL
        )
        if diagram_match:
            diagram_block = diagram_match.group(1)
            slide.diagram_spec = parse_diagram_block(diagram_block)
            slide.slide_type = 'text_diagram'

        # Extract notes
        notes_match = re.search(r'\*\*Notes\*\*:\s*(.+?)(?:\n\*\*|\Z)', slide_content, re.DOTALL)
        if notes_match:
            slide.notes = notes_match.group(1).strip()

        # Determine slide type
        if slide.number == 1:
            slide.slide_type = 'title'
        elif slide.slide_type not in ('text_image', 'text_diagram', 'section'):
            slide.slide_type = 'text_only'

        slides.append(slide)

    return slides


def emu_to_inches(emu) -> float:
    """Convert EMUs to inches."""
    return emu / 914400


def get_image_dimensions(image_path: Path) -> Tuple[int, int]:
    """Get image dimensions in pixels."""
    with Image.open(image_path) as img:
        return img.size


def find_placeholder(slide, placeholder_type=None, idx=None):
    """Find a placeholder by type or index."""
    for shape in slide.placeholders:
        if idx is not None and shape.placeholder_format.idx == idx:
            return shape
        if placeholder_type is not None:
            if shape.placeholder_format.type == placeholder_type:
                return shape
    return None


def find_body_placeholder(slide):
    """Find the body/content placeholder."""
    for shape in slide.placeholders:
        ph_type = shape.placeholder_format.type
        if ph_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
            return shape
    return None


def fill_title(slide, title_text: str):
    """Fill the title placeholder."""
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = title_text


def fill_body(slide, content: List[str], resize_width: float = None):
    """Fill the body placeholder with bullet content."""
    body_shape = find_body_placeholder(slide)
    if not body_shape:
        return

    # IMPORTANT: Save ALL original position/size before any modifications
    # Setting only some properties creates incomplete XML (python-pptx bug)
    original_left = body_shape.left
    original_top = body_shape.top
    original_width = body_shape.width
    original_height = body_shape.height

    # Resize width if needed - must set ALL FOUR properties together
    if resize_width:
        body_shape.left = original_left
        body_shape.top = original_top
        body_shape.width = Inches(resize_width)
        body_shape.height = original_height

    tf = body_shape.text_frame
    tf.clear()

    # Restore all properties after clear (clear can reset auto-fit)
    body_shape.left = original_left
    body_shape.top = original_top
    body_shape.width = Inches(resize_width) if resize_width else original_width
    body_shape.height = original_height

    for i, item in enumerate(content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Skip figure references
        if item.startswith('[Figure:') or item.startswith('**Figure**'):
            continue

        p.text = item
        p.level = 0


def insert_image(slide, image_path: Path, left: float, top: float,
                 max_width: float = None, max_height: float = None):
    """Insert an image as a shape with proper sizing."""

    if not image_path.exists():
        print(f"  Warning: Image not found: {image_path}")
        return None

    # Get image dimensions
    img_width_px, img_height_px = get_image_dimensions(image_path)
    aspect_ratio = img_width_px / img_height_px

    # Calculate size maintaining aspect ratio
    if max_width and max_height:
        # Fit within bounds
        if aspect_ratio > max_width / max_height:
            width = max_width
            height = width / aspect_ratio
        else:
            height = max_height
            width = height * aspect_ratio
    elif max_width:
        width = max_width
        height = width / aspect_ratio
    elif max_height:
        height = max_height
        width = height * aspect_ratio
    else:
        # Default size
        width = 5.0
        height = width / aspect_ratio

    # Add the picture
    pic = slide.shapes.add_picture(
        str(image_path),
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height)
    )

    return pic


def generate_presentation(
    template_path: Path,
    slides_data: List[SlideContent],
    figures_dir: Path,
    output_path: Path,
    layout_config: LayoutConfig
):
    """Generate the presentation with text and images."""

    prs = Presentation(str(template_path))

    # Get slide dimensions
    slide_width = emu_to_inches(prs.slide_width)
    slide_height = emu_to_inches(prs.slide_height)

    print(f"Slide dimensions: {slide_width:.2f}\" x {slide_height:.2f}\"")
    print(f"Processing {len(slides_data)} slides...")
    print()

    # We need to match slides to the existing presentation structure
    # If we have a working file with correct number of slides, use it

    if len(prs.slides) != len(slides_data):
        print(f"Warning: Template has {len(prs.slides)} slides, but content has {len(slides_data)} slides")
        print("Processing only matching slides...")

    image_count = 0
    diagram_count = 0

    for i, slide_content in enumerate(slides_data):
        if i >= len(prs.slides):
            break

        slide = prs.slides[i]

        # Fill title
        if slide_content.title:
            fill_title(slide, slide_content.title)

        # Handle based on slide type
        if slide_content.slide_type == 'text_image' and slide_content.figure_path:
            # Fill content (excluding figure references) with resized width
            clean_content = [
                c for c in slide_content.content
                if not c.startswith('[Figure:') and
                   not c.startswith('**Figure**') and
                   'status-tab' not in c.lower() and
                   '.png]' not in c
            ]
            # Pass resize_width to fill_body - it will handle width AND preserve height
            fill_body(slide, clean_content, resize_width=layout_config.text_width)

            # Resolve figure path
            fig_path = slide_content.figure_path
            if fig_path.startswith('../'):
                fig_path = fig_path[3:]  # Remove ../

            image_path = figures_dir.parent / fig_path
            if not image_path.exists():
                image_path = figures_dir / Path(fig_path).name

            if image_path.exists():
                # Calculate available space for image
                available_height = slide_height - layout_config.image_top - 0.5

                pic = insert_image(
                    slide,
                    image_path,
                    left=layout_config.image_left,
                    top=layout_config.image_top,
                    max_width=layout_config.image_width,
                    max_height=available_height
                )

                if pic:
                    image_count += 1
                    actual_left = emu_to_inches(pic.left)
                    actual_top = emu_to_inches(pic.top)
                    actual_width = emu_to_inches(pic.width)
                    actual_height = emu_to_inches(pic.height)
                    print(f"Slide {i+1}: {slide_content.title[:40]}...")
                    print(f"  Image: {image_path.name}")
                    print(f"  Position: ({actual_left:.2f}\", {actual_top:.2f}\") "
                          f"Size: {actual_width:.2f}\" x {actual_height:.2f}\"")
            else:
                print(f"Slide {i+1}: Image not found: {fig_path}")

        elif slide_content.slide_type == 'text_diagram' and slide_content.diagram_spec:
            # Fill content with resized width (same as images)
            fill_body(slide, slide_content.content, resize_width=layout_config.text_width)

            # Configure diagram position
            diagram_config = DiagramConfig(
                left=layout_config.image_left,
                top=layout_config.image_top,
                width=layout_config.image_width,
                height=slide_height - layout_config.image_top - 0.5
            )

            # Render the diagram
            shapes = render_diagram(slide, slide_content.diagram_spec, diagram_config)

            if shapes:
                diagram_count += 1
                diagram_type = slide_content.diagram_spec.get('type', 'unknown')
                node_count = len(slide_content.diagram_spec.get('nodes', []))
                print(f"Slide {i+1}: {slide_content.title[:40]}...")
                print(f"  Diagram: {diagram_type} with {node_count} nodes, {len(shapes)} shapes")

        elif slide_content.slide_type in ('text_only', 'section'):
            # Full-width text
            fill_body(slide, slide_content.content)

        elif slide_content.slide_type == 'title':
            # Title slide - just title and maybe subtitle
            fill_body(slide, slide_content.content)

    # Save
    prs.save(str(output_path))

    print()
    print(f"Generated {len(slides_data)} slides with {image_count} images and {diagram_count} diagrams")
    print(f"Saved to: {output_path}")


def main():
    parser = argparse.ArgumentParser(
        description="Generate PowerPoint from PRESENTATION.md with images",
        formatter_class=argparse.RawDescriptionHelpFormatter
    )
    parser.add_argument(
        "--template", "-t",
        required=True,
        help="PowerPoint template file (with slides already arranged)"
    )
    parser.add_argument(
        "--source", "-s",
        required=True,
        help="PRESENTATION.md source file"
    )
    parser.add_argument(
        "--figures", "-f",
        required=True,
        help="Figures directory"
    )
    parser.add_argument(
        "--output", "-o",
        required=True,
        help="Output PowerPoint file"
    )
    parser.add_argument(
        "--text-width",
        type=float,
        default=6.5,
        help="Width of text area for image slides (default: 6.5)"
    )
    parser.add_argument(
        "--image-left",
        type=float,
        default=7.0,
        help="Left position of images (default: 7.0)"
    )
    parser.add_argument(
        "--image-width",
        type=float,
        default=5.5,
        help="Max width of images (default: 5.5)"
    )
    parser.add_argument(
        "--image-top",
        type=float,
        default=1.5,
        help="Top position of images (default: 1.5)"
    )

    args = parser.parse_args()

    template_path = Path(args.template)
    source_path = Path(args.source)
    figures_dir = Path(args.figures)
    output_path = Path(args.output)

    if not template_path.exists():
        print(f"Error: Template not found: {template_path}")
        sys.exit(1)

    if not source_path.exists():
        print(f"Error: Source not found: {source_path}")
        sys.exit(1)

    # Parse PRESENTATION.md
    print(f"Parsing: {source_path}")
    slides_data = parse_presentation_md(source_path)
    print(f"Found {len(slides_data)} slides")

    # Count slides with figures
    image_slides = [s for s in slides_data if s.figure_path]
    print(f"Slides with figures: {len(image_slides)}")
    for s in image_slides:
        print(f"  Slide {s.number}: {s.figure_path}")

    # Count slides with diagrams
    diagram_slides = [s for s in slides_data if s.diagram_spec]
    print(f"Slides with diagrams: {len(diagram_slides)}")
    for s in diagram_slides:
        dtype = s.diagram_spec.get('type', 'unknown')
        nodes = len(s.diagram_spec.get('nodes', []))
        print(f"  Slide {s.number}: {dtype} diagram ({nodes} nodes)")
    print()

    # Configure layouts
    layout_config = LayoutConfig(
        text_width=args.text_width,
        image_left=args.image_left,
        image_width=args.image_width,
        image_top=args.image_top
    )

    # Generate
    generate_presentation(
        template_path,
        slides_data,
        figures_dir,
        output_path,
        layout_config
    )


if __name__ == "__main__":
    main()
