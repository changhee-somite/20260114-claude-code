#!/usr/bin/env python3
"""
Table Renderer: Create tables in PowerPoint presentations.

This module renders tables from markdown or structured specifications
using python-pptx. Supports styling, headers, and flexible positioning.

Features:
- Parse markdown tables or JSON specifications
- Customizable styling (colors, fonts, borders)
- Header row formatting
- Automatic column width calculation
- Position offset bug handling

Usage:
    from table_renderer import render_table, parse_markdown_table

    render_table(slide, table_spec, position_config)

Author: Claude Code
Date: 2026-01-18
"""

from dataclasses import dataclass
from typing import List, Dict, Optional, Tuple

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.table import _Cell


@dataclass
class TableConfig:
    """Configuration for table positioning."""
    left: float = 7.0       # Left position in inches
    top: float = 1.5        # Top position in inches
    width: float = 5.0      # Total width in inches
    height: float = None    # Height (auto-calculated if None)
    row_height: float = 0.4 # Default row height in inches


@dataclass
class TableStyle:
    """Styling options for tables."""
    header_fill: str = "#4472C4"      # Header background color (blue)
    header_text: str = "#FFFFFF"       # Header text color (white)
    row_fill: str = "#FFFFFF"          # Row background color
    alt_row_fill: str = "#F2F2F2"      # Alternating row background
    text_color: str = "#000000"        # Text color
    border_color: str = "#A6A6A6"      # Border color
    font_size: int = 10                # Font size in points
    header_font_size: int = 11         # Header font size
    header_bold: bool = True           # Bold header text
    use_alt_rows: bool = True          # Use alternating row colors


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color string to RGBColor."""
    hex_color = hex_color.lstrip('#')
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


def parse_markdown_table(md_text: str) -> Dict:
    """
    Parse a markdown table into a structured specification.

    Args:
        md_text: Markdown table text

    Returns:
        Dict with 'headers' and 'rows' keys
    """
    lines = [l.strip() for l in md_text.strip().split('\n') if l.strip()]

    if not lines:
        return {"headers": [], "rows": []}

    result = {"headers": [], "rows": []}

    for i, line in enumerate(lines):
        # Skip separator lines (|---|---|)
        if set(line.replace('|', '').replace('-', '').replace(':', '').strip()) == set():
            continue
        if '---' in line or '|--' in line:
            continue

        # Parse cells
        cells = [c.strip() for c in line.split('|')]
        # Remove empty strings from start/end (from leading/trailing |)
        cells = [c for c in cells if c or cells.index(c) not in (0, len(cells)-1)]
        cells = [c for c in cells if c]  # Remove empty cells

        if not cells:
            continue

        if i == 0:
            result["headers"] = cells
        else:
            result["rows"].append(cells)

    return result


def calculate_column_widths(headers: List[str], rows: List[List[str]],
                            total_width: float) -> List[float]:
    """
    Calculate column widths based on content.

    Uses a weighted approach based on character count.
    """
    if not headers:
        return []

    num_cols = len(headers)

    # Calculate max width per column (in characters)
    max_chars = [len(h) for h in headers]

    for row in rows:
        for i, cell in enumerate(row):
            if i < len(max_chars):
                max_chars[i] = max(max_chars[i], len(cell))

    # Calculate proportional widths
    total_chars = sum(max_chars) or 1
    widths = [(chars / total_chars) * total_width for chars in max_chars]

    # Ensure minimum width
    min_width = 0.5
    for i, w in enumerate(widths):
        if w < min_width:
            widths[i] = min_width

    return widths


def apply_formatted_text_to_cell(paragraph, text: str, text_color: str,
                                  font_size: int, base_bold: bool = False):
    """
    Apply text to a paragraph, parsing **bold** markdown patterns.

    Args:
        paragraph: python-pptx paragraph object
        text: Text that may contain **bold** patterns
        text_color: Hex color for text
        font_size: Font size in points
        base_bold: Whether all text should be bold (for headers)
    """
    import re

    # Clear existing content
    paragraph.clear()

    # Pattern to find **bold** text
    bold_pattern = re.compile(r'\*\*(.+?)\*\*')

    last_end = 0
    for match in bold_pattern.finditer(text):
        # Add non-bold text before this match
        if match.start() > last_end:
            run = paragraph.add_run()
            run.text = text[last_end:match.start()]
            run.font.size = Pt(font_size)
            run.font.color.rgb = hex_to_rgb(text_color)
            run.font.bold = base_bold

        # Add bold text
        run = paragraph.add_run()
        run.text = match.group(1)
        run.font.size = Pt(font_size)
        run.font.color.rgb = hex_to_rgb(text_color)
        run.font.bold = True  # Always bold for **text**

        last_end = match.end()

    # Add any remaining non-bold text after the last match
    # OR if no matches were found (last_end == 0), add the whole text
    if last_end < len(text):
        run = paragraph.add_run()
        run.text = text[last_end:]
        run.font.size = Pt(font_size)
        run.font.color.rgb = hex_to_rgb(text_color)
        run.font.bold = base_bold


def style_cell(cell: _Cell, text: str, fill_color: str, text_color: str,
               font_size: int, bold: bool = False, align: PP_ALIGN = PP_ALIGN.LEFT):
    """Apply styling to a table cell with markdown support."""
    # Fill color
    cell.fill.solid()
    cell.fill.fore_color.rgb = hex_to_rgb(fill_color)

    # Text formatting with markdown parsing - use only the first paragraph
    # (cells may have multiple default paragraphs which would cause duplication)
    if cell.text_frame.paragraphs:
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.alignment = align
        apply_formatted_text_to_cell(paragraph, text, text_color, font_size, bold)

    # Vertical alignment
    cell.text_frame.word_wrap = True
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def render_table(slide, spec: Dict, config: TableConfig = None,
                 style: TableStyle = None) -> Optional[object]:
    """
    Render a table on a slide based on the specification.

    Args:
        slide: PowerPoint slide object
        spec: Table specification dict with 'headers' and 'rows'
        config: TableConfig for positioning
        style: TableStyle for appearance

    Returns:
        Table shape object or None if failed
    """
    if config is None:
        config = TableConfig()
    if style is None:
        style = TableStyle()

    headers = spec.get("headers", [])
    rows = spec.get("rows", [])

    if not headers:
        print("  Warning: No headers provided for table")
        return None

    num_rows = len(rows) + 1  # +1 for header row
    num_cols = len(headers)

    # Calculate height
    if config.height:
        height = config.height
    else:
        height = num_rows * config.row_height

    # Calculate column widths
    col_widths = calculate_column_widths(headers, rows, config.width)

    # CRITICAL: Use explicit position values to avoid python-pptx offset bug
    # All four properties must be set to generate complete XML
    left = Inches(config.left)
    top = Inches(config.top)
    width = Inches(config.width)
    height_emu = Inches(height)

    # Create table
    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        left, top, width, height_emu
    )
    table = table_shape.table

    # Set column widths
    for i, col_width in enumerate(col_widths):
        if i < len(table.columns):
            table.columns[i].width = Inches(col_width)

    # Style header row
    for col_idx, header_text in enumerate(headers):
        if col_idx < num_cols:
            cell = table.cell(0, col_idx)
            style_cell(
                cell, header_text,
                fill_color=style.header_fill,
                text_color=style.header_text,
                font_size=style.header_font_size,
                bold=style.header_bold,
                align=PP_ALIGN.CENTER
            )

    # Style data rows
    for row_idx, row_data in enumerate(rows):
        # Alternating row colors
        if style.use_alt_rows and row_idx % 2 == 1:
            fill = style.alt_row_fill
        else:
            fill = style.row_fill

        for col_idx, cell_text in enumerate(row_data):
            if col_idx < num_cols:
                cell = table.cell(row_idx + 1, col_idx)
                style_cell(
                    cell, cell_text,
                    fill_color=fill,
                    text_color=style.text_color,
                    font_size=style.font_size,
                    bold=False,
                    align=PP_ALIGN.LEFT
                )

    return table_shape


def parse_table_block(block: str) -> Dict:
    """
    Parse a ```table code block into a specification dict.

    Supports two formats:

    1. Markdown table:
        ```table
        | Header 1 | Header 2 |
        |----------|----------|
        | Cell 1   | Cell 2   |
        ```

    2. YAML-like specification:
        ```table
        headers:
          - Header 1
          - Header 2
        rows:
          - [Cell 1, Cell 2]
          - [Cell 3, Cell 4]
        style:
          header_fill: "#4472C4"
        ```
    """
    block = block.strip()

    # Detect format: markdown tables start with |
    if block.startswith('|'):
        return parse_markdown_table(block)

    # Parse YAML-like format
    spec = {
        "headers": [],
        "rows": [],
        "style": {}
    }

    lines = block.split('\n')
    current_section = None
    current_row = None

    for line in lines:
        stripped = line.strip()

        if not stripped or stripped.startswith('#'):
            continue

        if stripped == 'headers:':
            current_section = 'headers'
        elif stripped == 'rows:':
            current_section = 'rows'
        elif stripped == 'style:':
            current_section = 'style'
        elif stripped.startswith('- [') and current_section == 'rows':
            # Parse row as list: - [Cell 1, Cell 2]
            row_content = stripped[3:-1]  # Remove "- [" and "]"
            cells = [c.strip().strip('"\'') for c in row_content.split(',')]
            spec["rows"].append(cells)
        elif stripped.startswith('- ') and current_section == 'headers':
            spec["headers"].append(stripped[2:].strip())
        elif ':' in stripped and current_section == 'style':
            key, value = stripped.split(':', 1)
            spec["style"][key.strip()] = value.strip().strip('"\'')
        elif stripped.startswith('position:'):
            spec["position"] = stripped.split(':', 1)[1].strip()
        elif stripped.startswith('width:'):
            spec["width"] = float(stripped.split(':', 1)[1].strip())
        elif stripped.startswith('top:'):
            spec["top"] = float(stripped.split(':', 1)[1].strip())
        elif stripped.startswith('left:'):
            spec["left"] = float(stripped.split(':', 1)[1].strip())

    return spec


# For convenience, export key functions
__all__ = ['render_table', 'parse_markdown_table', 'parse_table_block',
           'TableConfig', 'TableStyle']
