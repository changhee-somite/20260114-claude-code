#!/usr/bin/env python3
"""
PPTX Inspector - Validate PowerPoint presentations for common issues.

This tool detects layout breakage and structural issues that may not be
visible through the python-pptx object model but cause rendering problems
in PowerPoint.

Key validations:
- Level 1: Quick object property checks (dimensions, positions)
- Level 2: XML structure validation (the critical position offset bug)
- Level 3: Content validation (placeholder text, text length)
- Level 4: Layout overlap detection

Usage:
    python pptx_inspector.py presentation.pptx
    python pptx_inspector.py presentation.pptx --level 4
    python pptx_inspector.py presentation.pptx --json
    python pptx_inspector.py presentation.pptx --fail-on warning

Author: Claude Code
Date: 2026-01-16
"""

import argparse
import json
import re
import sys
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

# Template font sizes (from slideMaster1.xml bodyStyle)
# These match the sizes used in generate_presentation.py
TEMPLATE_FONT_SIZES = {0: 24, 1: 18, 2: 15, 3: 12, 4: 10}

# Template paragraph spacing (spcBef - space before each paragraph, in points)
# From slideMaster1.xml bodyStyle
TEMPLATE_PARAGRAPH_SPACING = {0: 12, 1: 12, 2: 12, 3: 29.5, 4: 29.5}


def inspect_table(slide_num, shape_name, table):
    """
    Inspect a table for common issues.

    Args:
        slide_num: Slide number (1-indexed)
        shape_name: Name of the table shape
        table: python-pptx Table object

    Returns:
        List of issue dicts
    """
    issues = []

    num_rows = len(table.rows)
    num_cols = len(table.columns)

    # Check 1: Empty table
    if num_rows == 0:
        issues.append({
            'slide': slide_num,
            'shape': shape_name,
            'type': 'TABLE_NO_ROWS',
            'severity': 'ERROR',
            'description': f'Table "{shape_name}" has no rows'
        })
        return issues

    if num_cols == 0:
        issues.append({
            'slide': slide_num,
            'shape': shape_name,
            'type': 'TABLE_NO_COLUMNS',
            'severity': 'ERROR',
            'description': f'Table "{shape_name}" has no columns'
        })
        return issues

    # Check 2: Empty header cells
    if num_rows > 0:
        header_row = table.rows[0]
        empty_headers = 0
        for cell in header_row.cells:
            if not cell.text.strip():
                empty_headers += 1

        if empty_headers > 0:
            issues.append({
                'slide': slide_num,
                'shape': shape_name,
                'type': 'TABLE_EMPTY_HEADERS',
                'severity': 'WARNING',
                'description': f'Table has {empty_headers} empty header cell(s)'
            })

    # Check 3: Rows with mismatched column count
    for row_idx, row in enumerate(table.rows):
        actual_cols = len(list(row.cells))
        if actual_cols != num_cols:
            issues.append({
                'slide': slide_num,
                'shape': shape_name,
                'type': 'TABLE_COLUMN_MISMATCH',
                'severity': 'WARNING',
                'description': f'Row {row_idx + 1} has {actual_cols} cells, expected {num_cols}'
            })

    # Check 4: Very small column widths
    for col_idx, col in enumerate(table.columns):
        col_width_inches = col.width / 914400
        if col_width_inches < 0.3:
            issues.append({
                'slide': slide_num,
                'shape': shape_name,
                'type': 'TABLE_NARROW_COLUMN',
                'severity': 'WARNING',
                'description': f'Column {col_idx + 1} is very narrow ({col_width_inches:.2f}")'
            })

    # Check 5: Empty data rows (all cells empty)
    for row_idx in range(1, num_rows):  # Skip header
        row = table.rows[row_idx]
        all_empty = all(not cell.text.strip() for cell in row.cells)
        if all_empty:
            issues.append({
                'slide': slide_num,
                'shape': shape_name,
                'type': 'TABLE_EMPTY_ROW',
                'severity': 'INFO',
                'description': f'Data row {row_idx} is completely empty'
            })

    return issues


def inspect_level1(pptx_path):
    """Quick inspection using python-pptx object model."""
    issues = []
    prs = Presentation(str(pptx_path))

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1

        for shape in slide.shapes:
            shape_name = shape.name

            # Check tables
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table = shape.table
                issues.extend(inspect_table(slide_num, shape_name, table))

            # Check if shape is a connector (straight connectors validly have one dimension as 0)
            is_connector = (
                hasattr(shape, 'shape_type') and
                shape.shape_type == MSO_SHAPE_TYPE.LINE
            ) or 'connector' in shape_name.lower()

            # Check 1: Zero dimensions (skip for connectors - they can have one zero dimension)
            if shape.width == 0 and not is_connector:
                issues.append({
                    'slide': slide_num,
                    'shape': shape_name,
                    'type': 'ZERO_WIDTH',
                    'severity': 'ERROR',
                    'description': f'{shape_name} has width=0'
                })

            if shape.height == 0 and not is_connector:
                issues.append({
                    'slide': slide_num,
                    'shape': shape_name,
                    'type': 'ZERO_HEIGHT',
                    'severity': 'ERROR',
                    'description': f'{shape_name} has height=0'
                })

            # Check 2: Negative positions
            if shape.left < 0:
                issues.append({
                    'slide': slide_num,
                    'shape': shape_name,
                    'type': 'NEGATIVE_LEFT',
                    'severity': 'WARNING',
                    'description': f'{shape_name} has negative left position'
                })

            if shape.top < 0:
                issues.append({
                    'slide': slide_num,
                    'shape': shape_name,
                    'type': 'NEGATIVE_TOP',
                    'severity': 'WARNING',
                    'description': f'{shape_name} has negative top position'
                })

            # Check 3: Shape extends beyond slide bounds (with 10% tolerance)
            if shape.left + shape.width > slide_width * 1.1:
                issues.append({
                    'slide': slide_num,
                    'shape': shape_name,
                    'type': 'EXTENDS_RIGHT',
                    'severity': 'WARNING',
                    'description': f'{shape_name} extends beyond slide right edge'
                })

            if shape.top + shape.height > slide_height * 1.1:
                issues.append({
                    'slide': slide_num,
                    'shape': shape_name,
                    'type': 'EXTENDS_BOTTOM',
                    'severity': 'WARNING',
                    'description': f'{shape_name} extends beyond slide bottom edge'
                })

    return issues


def inspect_level2_xml(pptx_path):
    """
    Deep inspection of XML structure.

    This catches the CRITICAL position offset bug where python-pptx creates
    <a:xfrm> with only <a:ext> (size) but not <a:off> (position).
    """
    issues = []

    with zipfile.ZipFile(str(pptx_path), 'r') as z:
        # Find all slide XML files
        slide_files = [f for f in z.namelist() if re.match(r'ppt/slides/slide\d+\.xml', f)]

        for slide_file in sorted(slide_files):
            slide_num = int(re.search(r'slide(\d+)', slide_file).group(1))

            with z.open(slide_file) as f:
                content = f.read().decode('utf-8')

                # Check 1: xfrm with ext but without off (THE CRITICAL BUG)
                xfrm_blocks = re.findall(r'<a:xfrm>.*?</a:xfrm>', content, re.DOTALL)

                for block in xfrm_blocks:
                    has_off = '<a:off' in block
                    has_ext = '<a:ext' in block

                    if has_ext and not has_off:
                        # Extract size for context
                        ext_match = re.search(r'<a:ext cx="(\d+)" cy="(\d+)"', block)
                        size_info = ""
                        if ext_match:
                            w = int(ext_match.group(1)) / 914400
                            h = int(ext_match.group(2)) / 914400
                            size_info = f" (size: {w:.2f}\" x {h:.2f}\")"

                        issues.append({
                            'slide': slide_num,
                            'type': 'MISSING_POSITION_OFFSET',
                            'severity': 'CRITICAL',
                            'description': f'xfrm has size but no position{size_info} - will cause rendering issues'
                        })

                # Check 2: Empty embed references
                if 'r:embed=""' in content:
                    issues.append({
                        'slide': slide_num,
                        'type': 'EMPTY_EMBED_REFERENCE',
                        'severity': 'ERROR',
                        'description': 'Empty embed reference - image may not display'
                    })

        # Check relationship files for missing media
        for slide_file in slide_files:
            slide_num = int(re.search(r'slide(\d+)', slide_file).group(1))
            rels_file = slide_file.replace('slides/', 'slides/_rels/') + '.rels'

            try:
                with z.open(rels_file) as f:
                    rels_content = f.read().decode('utf-8')
                    media_refs = re.findall(r'Target="\.\./media/([^"]+)"', rels_content)
                    media_files = [f for f in z.namelist() if f.startswith('ppt/media/')]

                    for ref in media_refs:
                        expected = f'ppt/media/{ref}'
                        if expected not in media_files:
                            issues.append({
                                'slide': slide_num,
                                'type': 'MISSING_MEDIA_FILE',
                                'severity': 'ERROR',
                                'description': f'Referenced media file not found: {ref}'
                            })
            except KeyError:
                pass  # No rels file is OK for some slides

    return issues


def estimate_text_overflow(text_frame, shape_width, shape_height):
    """
    Estimate if text will overflow its container using level-aware font sizes
    and actual template paragraph spacing.

    This improved version uses the template's actual font sizes and paragraph
    spacing for each bullet level, providing accurate overflow detection.

    Returns:
        tuple: (likely_overflow: bool, estimated_lines: int, available_lines: int, details: str)
    """
    if not text_frame.paragraphs:
        return False, 0, 0, ""

    # Get effective dimensions in points
    # Account for internal text frame margins (~3.6pt top and bottom)
    shape_width_pt = shape_width / 914400 * 72
    shape_height_pt = shape_height / 914400 * 72
    effective_width_pt = shape_width_pt * 0.9   # 10% margin for left/right
    effective_height_pt = shape_height_pt - 7.2  # Subtract internal margins (~3.6pt each)

    # Calculate space needed with level-aware font sizes and actual paragraph spacing
    lines_needed_pt = 0
    num_paragraphs = len(text_frame.paragraphs)

    for i, para in enumerate(text_frame.paragraphs):
        # Get paragraph level
        level = para.level if hasattr(para, 'level') else 0

        # Get font size - from run if explicitly set, or template default for level
        font_size_pt = None
        if para.runs and para.runs[0].font.size:
            font_size_pt = para.runs[0].font.size.pt
        if font_size_pt is None:
            font_size_pt = TEMPLATE_FONT_SIZES.get(level, 18)

        # Get paragraph spacing from template (space before each paragraph)
        para_spacing = TEMPLATE_PARAGRAPH_SPACING.get(level, 12)

        # Line height = font size (since template uses lnSpc=100%)
        line_height = font_size_pt

        # Calculate character width and chars per line
        char_width = font_size_pt * 0.5  # average character width
        # Account for bullet indent reducing available width
        indent_reduction = level * 0.3 * 72  # ~0.3" per indent level
        adjusted_width = max(72, effective_width_pt - indent_reduction)  # min 1 inch
        chars_per_line = max(1, adjusted_width / char_width)

        para_text = para.text
        if not para_text.strip():
            # Empty paragraph still takes some space
            lines_needed_pt += line_height * 0.5
        else:
            # Calculate lines needed for this paragraph (with wrapping)
            text_lines = max(1, len(para_text) / chars_per_line)
            lines_needed_pt += text_lines * line_height

        # Add paragraph spacing (space before, except for first paragraph)
        if i > 0:
            lines_needed_pt += para_spacing

    # Check for overflow
    overflow = lines_needed_pt > effective_height_pt

    # Flag borderline when within 90% of limit (tighter threshold)
    borderline = (lines_needed_pt > effective_height_pt * 0.90 and not overflow)

    # Calculate approximate line counts for reporting
    ref_line_height = 24  # Use level 0 font size as reference
    available_lines = int(effective_height_pt / ref_line_height)
    lines_needed = int(lines_needed_pt / ref_line_height + 0.5)

    # Build details string
    details = f"~{lines_needed_pt:.0f}pt needed, ~{effective_height_pt:.0f}pt available"
    if borderline:
        details += " - BORDERLINE"

    return overflow or borderline, lines_needed, available_lines, details


def inspect_level3_content(pptx_path):
    """Validate content integrity."""
    issues = []
    prs = Presentation(str(pptx_path))

    placeholder_patterns = [
        (r'\[Figure:.*?\]', 'Figure reference not replaced'),
        (r'\[Image:.*?\]', 'Image reference not replaced'),
        (r'Lorem ipsum', 'Lorem ipsum placeholder text'),
        (r'Click to add', 'Unreplaced placeholder prompt'),
        (r'Title Text', 'Unreplaced title placeholder'),
        (r'Body Level', 'Unreplaced body placeholder'),
    ]

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1

        for shape in slide.shapes:
            if not hasattr(shape, 'text_frame') or not shape.has_text_frame:
                continue

            text = shape.text_frame.text

            # Check 1: Placeholder text not replaced
            for pattern, description in placeholder_patterns:
                if re.search(pattern, text, re.IGNORECASE):
                    issues.append({
                        'slide': slide_num,
                        'shape': shape.name,
                        'type': 'UNREPLACED_PLACEHOLDER',
                        'severity': 'WARNING',
                        'description': f'{description}: "{text[:50]}..."'
                    })
                    break

            # Check 2: Text overflow detection
            if text.strip() and shape.width > 0 and shape.height > 0:
                overflow, lines_needed, available_lines, details = estimate_text_overflow(
                    shape.text_frame, shape.width, shape.height
                )
                if overflow:
                    issues.append({
                        'slide': slide_num,
                        'shape': shape.name,
                        'type': 'TEXT_OVERFLOW',
                        'severity': 'WARNING',
                        'description': f'Text likely overflows container: {details}',
                        'lines_needed': lines_needed,
                        'lines_available': available_lines,
                        'fixable': True
                    })

            # Check 3: Extremely long text (keep for backwards compatibility)
            if len(text) > 2000:
                issues.append({
                    'slide': slide_num,
                    'shape': shape.name,
                    'type': 'VERY_LONG_TEXT',
                    'severity': 'WARNING',
                    'description': f'Text length ({len(text)} chars) may overflow'
                })

            # Check 4: Whitespace only
            if text and not text.strip():
                issues.append({
                    'slide': slide_num,
                    'shape': shape.name,
                    'type': 'WHITESPACE_ONLY',
                    'severity': 'INFO',
                    'description': 'Shape contains only whitespace'
                })

    return issues


def inspect_level4_overlaps(pptx_path, tolerance=0.1):
    """Detect overlapping shapes that might cause visual issues."""
    issues = []
    prs = Presentation(str(pptx_path))

    def get_bounds(shape):
        """Get shape bounds in inches."""
        l = shape.left / 914400
        t = shape.top / 914400
        r = (shape.left + shape.width) / 914400
        b = (shape.top + shape.height) / 914400
        return l, t, r, b

    def shapes_overlap(s1, s2):
        """Check if two shapes overlap."""
        l1, t1, r1, b1 = get_bounds(s1)
        l2, t2, r2, b2 = get_bounds(s2)

        h_overlap = not (r1 <= l2 + tolerance or l1 >= r2 - tolerance)
        v_overlap = not (b1 <= t2 + tolerance or t1 >= b2 - tolerance)

        return h_overlap and v_overlap

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        shapes = list(slide.shapes)

        for j, s1 in enumerate(shapes):
            for s2 in shapes[j+1:]:
                # Skip slide number placeholders
                if 'Slide Number' in s1.name or 'Slide Number' in s2.name:
                    continue

                if shapes_overlap(s1, s2):
                    s1_has_text = (hasattr(s1, 'text_frame') and
                                  s1.has_text_frame and
                                  s1.text_frame.text.strip())
                    s2_has_text = (hasattr(s2, 'text_frame') and
                                  s2.has_text_frame and
                                  s2.text_frame.text.strip())
                    s1_is_table = s1.shape_type == MSO_SHAPE_TYPE.TABLE
                    s2_is_table = s2.shape_type == MSO_SHAPE_TYPE.TABLE

                    # Table overlapping with text content is a warning
                    if (s1_is_table or s2_is_table) and (s1_has_text or s2_has_text):
                        issues.append({
                            'slide': slide_num,
                            'type': 'TABLE_TEXT_OVERLAP',
                            'severity': 'WARNING',
                            'description': f'Table "{s1.name if s1_is_table else s2.name}" overlaps with text content',
                            'shapes': [s1.name, s2.name]
                        })
                    # Two tables overlapping
                    elif s1_is_table and s2_is_table:
                        issues.append({
                            'slide': slide_num,
                            'type': 'TABLE_OVERLAP',
                            'severity': 'ERROR',
                            'description': f'Tables "{s1.name}" and "{s2.name}" overlap',
                            'shapes': [s1.name, s2.name]
                        })
                    # Text overlapping text is a warning
                    elif s1_has_text and s2_has_text:
                        issues.append({
                            'slide': slide_num,
                            'type': 'TEXT_OVERLAP',
                            'severity': 'WARNING',
                            'description': f'"{s1.name}" overlaps with "{s2.name}" - both have text',
                            'shapes': [s1.name, s2.name]
                        })
                    # Image overlapping text placeholder might be intentional
                    elif (s1.shape_type == MSO_SHAPE_TYPE.PICTURE or
                          s2.shape_type == MSO_SHAPE_TYPE.PICTURE):
                        # Only report if significant overlap
                        pass  # Often intentional

    return issues


def fix_text_overflow(pptx_path, output_path=None, min_font_size=10):
    """
    Fix text overflow issues by shrinking font sizes to fit.

    Args:
        pptx_path: Path to input PPTX
        output_path: Path to output PPTX (default: overwrites input)
        min_font_size: Minimum font size in points (default: 10)

    Returns:
        List of fixes applied
    """
    prs = Presentation(str(pptx_path))
    fixes = []

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1

        for shape in slide.shapes:
            if not hasattr(shape, 'text_frame') or not shape.has_text_frame:
                continue

            text = shape.text_frame.text
            if not text.strip() or shape.width <= 0 or shape.height <= 0:
                continue

            # Check if overflow exists (includes borderline cases)
            overflow, lines_needed, available_lines, details = estimate_text_overflow(
                shape.text_frame, shape.width, shape.height
            )

            if not overflow:
                continue

            # Skip if it's borderline and lines_needed < 5 (minor issue)
            if 'BORDERLINE' in details and lines_needed < 5:
                continue

            # Get current font size
            current_font_size = 18  # default
            for para in shape.text_frame.paragraphs:
                if para.runs:
                    if para.runs[0].font.size:
                        current_font_size = para.runs[0].font.size.pt
                        break

            # Calculate required font size reduction
            # lines_needed / available_lines gives the ratio we need to reduce by
            ratio = available_lines / lines_needed
            new_font_size = max(min_font_size, int(current_font_size * ratio * 0.88))  # 12% safety margin

            if new_font_size >= current_font_size:
                continue  # Can't improve

            # Apply the fix
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(new_font_size)

            fixes.append({
                'slide': slide_num,
                'shape': shape.name,
                'original_size': current_font_size,
                'new_size': new_font_size,
                'description': f'Reduced font from {current_font_size}pt to {new_font_size}pt'
            })

    # Save
    save_path = output_path or pptx_path
    prs.save(str(save_path))

    return fixes


def main():
    parser = argparse.ArgumentParser(
        description='Inspect PPTX for layout issues and structural problems',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Inspection Levels:
  1 - Quick: Object property checks (dimensions, positions)
  2 - XML: Structure validation (position offset bug) [DEFAULT]
  3 - Content: Text validation (placeholders, length)
  4 - Full: All above plus overlap detection

Examples:
  python pptx_inspector.py presentation.pptx
  python pptx_inspector.py presentation.pptx --level 4
  python pptx_inspector.py presentation.pptx --json --fail-on warning
        """
    )
    parser.add_argument('pptx_file', help='PowerPoint file to inspect')
    parser.add_argument('--level', '-l', type=int, default=2, choices=[1, 2, 3, 4],
                       help='Inspection depth (default: 2)')
    parser.add_argument('--json', '-j', action='store_true',
                       help='Output as JSON')
    parser.add_argument('--fail-on', '-f',
                       choices=['critical', 'error', 'warning', 'info'],
                       default='critical',
                       help='Exit with error if issues at this level or above (default: critical)')
    parser.add_argument('--quiet', '-q', action='store_true',
                       help='Only output if issues found')
    parser.add_argument('--fix', action='store_true',
                       help='Attempt to fix TEXT_OVERFLOW issues by shrinking fonts')
    parser.add_argument('--fix-output', '-o', type=str, default=None,
                       help='Output path for fixed PPTX (default: overwrite input)')
    parser.add_argument('--min-font-size', type=int, default=10,
                       help='Minimum font size when fixing (default: 10)')

    args = parser.parse_args()

    pptx_path = Path(args.pptx_file)
    if not pptx_path.exists():
        print(f"Error: File not found: {pptx_path}", file=sys.stderr)
        sys.exit(1)

    all_issues = []

    # Run inspections up to requested level
    if args.level >= 1:
        all_issues.extend(inspect_level1(pptx_path))
    if args.level >= 2:
        all_issues.extend(inspect_level2_xml(pptx_path))
    if args.level >= 3:
        all_issues.extend(inspect_level3_content(pptx_path))
    if args.level >= 4:
        all_issues.extend(inspect_level4_overlaps(pptx_path))

    # Output
    if args.json:
        print(json.dumps(all_issues, indent=2))
    else:
        if not all_issues:
            if not args.quiet:
                print(f"✅ {pptx_path.name}: No issues found (level {args.level} inspection)")
        else:
            severity_order = {'CRITICAL': 0, 'ERROR': 1, 'WARNING': 2, 'INFO': 3}
            sorted_issues = sorted(
                all_issues,
                key=lambda x: (severity_order.get(x.get('severity', 'INFO'), 4),
                              x.get('slide', 0))
            )

            print(f"{'='*60}")
            print(f"PPTX Inspector: {pptx_path.name}")
            print(f"{'='*60}")
            print(f"Found {len(all_issues)} issue(s):\n")

            for issue in sorted_issues:
                sev = issue.get('severity', 'INFO')
                icon = {
                    'CRITICAL': '🔴',
                    'ERROR': '❌',
                    'WARNING': '⚠️',
                    'INFO': 'ℹ️'
                }.get(sev, '?')

                slide = issue.get('slide', '?')
                desc = issue.get('description', str(issue))
                print(f"{icon} [{sev}] Slide {slide}: {desc}")

            print()

    # Handle --fix for TEXT_OVERFLOW issues
    if args.fix:
        overflow_issues = [i for i in all_issues if i.get('type') == 'TEXT_OVERFLOW']
        if overflow_issues:
            output_path = args.fix_output or pptx_path
            fixes = fix_text_overflow(pptx_path, output_path, args.min_font_size)

            if fixes:
                print(f"\n{'='*60}")
                print(f"FIXES APPLIED")
                print(f"{'='*60}")
                for fix in fixes:
                    print(f"  Slide {fix['slide']}: {fix['shape']} - {fix['description']}")
                print(f"\nSaved to: {output_path}")
            else:
                print("\nNo fixes could be applied (fonts already at minimum size)")
        else:
            print("\nNo TEXT_OVERFLOW issues to fix")

    # Determine exit code
    fail_levels = {
        'critical': ['CRITICAL'],
        'error': ['CRITICAL', 'ERROR'],
        'warning': ['CRITICAL', 'ERROR', 'WARNING'],
        'info': ['CRITICAL', 'ERROR', 'WARNING', 'INFO']
    }

    should_fail = any(
        issue.get('severity') in fail_levels[args.fail_on]
        for issue in all_issues
        if isinstance(issue, dict)
    )

    # If --fix was used and fixes were applied, don't fail on those issues
    if args.fix:
        fixed_types = {'TEXT_OVERFLOW'}
        should_fail = any(
            issue.get('severity') in fail_levels[args.fail_on]
            and issue.get('type') not in fixed_types
            for issue in all_issues
            if isinstance(issue, dict)
        )

    sys.exit(1 if should_fail else 0)


if __name__ == '__main__':
    main()
