#!/usr/bin/env python3
"""
Layout Discovery Module

Dynamically analyzes PowerPoint template layouts to determine their capabilities
and placeholder structure, enabling the presentation generator to work with
any template rather than relying on hardcoded layout indices.

Author: Claude Code
Date: 2026-01-30
"""

from dataclasses import dataclass, field
from typing import Dict, List, Optional, Tuple
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER


@dataclass
class PlaceholderInfo:
    """Information about a single placeholder in a layout."""
    idx: int                  # Placeholder index
    type: int                 # PP_PLACEHOLDER enum value
    type_name: str            # Human-readable type name
    left: float               # Position in inches
    top: float
    width: float
    height: float
    center_x: float           # Center position for arrangement detection
    center_y: float


@dataclass
class LayoutCapabilities:
    """Capabilities and placeholder indices for a single layout."""
    index: int                        # Layout index in template
    name: str                         # Layout display name

    # Capability flags
    has_title: bool = False
    has_subtitle: bool = False        # Body placeholder positioned as subtitle
    has_single_body: bool = False     # Single BODY placeholder (full-width text)
    has_single_object: bool = False   # Single OBJECT placeholder (full-width content)
    has_two_content: bool = False     # Two content placeholders

    # Content arrangement (for two-content layouts)
    content_arrangement: str = "none"  # "none", "side_by_side", "stacked"

    # Placeholder indices for each role
    title_idx: Optional[int] = None
    subtitle_idx: Optional[int] = None
    body_idx: Optional[int] = None

    # For side-by-side layouts
    left_content_idx: Optional[int] = None
    right_content_idx: Optional[int] = None

    # For stacked layouts
    top_content_idx: Optional[int] = None
    bottom_content_idx: Optional[int] = None

    # Whether placeholders have typed BODY vs OBJECT distinction
    has_typed_placeholders: bool = False

    # Store all placeholder info for fallback positioning
    _placeholders: Dict[int, PlaceholderInfo] = field(default_factory=dict)

    def get_placeholder_info(self, idx: int) -> Optional[PlaceholderInfo]:
        """Get placeholder info by index."""
        return self._placeholders.get(idx)


@dataclass
class LayoutMapping:
    """Mapping of slide types to discovered layouts."""
    title: Optional[LayoutCapabilities] = None
    section: Optional[LayoutCapabilities] = None
    text_only: Optional[LayoutCapabilities] = None
    text_image: Optional[LayoutCapabilities] = None
    text_image_top: Optional[LayoutCapabilities] = None
    text_diagram: Optional[LayoutCapabilities] = None
    text_table: Optional[LayoutCapabilities] = None
    text_table_full: Optional[LayoutCapabilities] = None  # Full-width table (no text)
    content_only: Optional[LayoutCapabilities] = None     # Full-width OBJECT placeholder

    # All discovered layouts for reference
    all_layouts: List[LayoutCapabilities] = field(default_factory=list)

    def get_layout_for_type(self, slide_type: str) -> Tuple[Optional[LayoutCapabilities], Dict]:
        """
        Get the best layout for a slide type, with fallback configuration.

        Returns:
            Tuple of (layout_caps, fallback_config) where fallback_config contains
            manual positioning overrides when using a non-ideal layout.
        """
        # Primary assignments
        primary = {
            'title': self.title,
            'section': self.section or self.title,
            'text_only': self.text_only,
            'text_image': self.text_image,
            'text_image_top': self.text_image_top,
            'text_diagram': self.text_diagram or self.text_image,
            'text_table': self.text_table or self.text_image,  # Split layout for tables with text
            'text_table_full': self.text_table_full or self.content_only or self.text_only,  # Full-width table
            'content_only': self.content_only or self.text_only,
        }

        layout = primary.get(slide_type)
        fallback_config = {}

        if layout is not None:
            return layout, fallback_config

        # Fallback chain when primary is None
        if slide_type == 'title':
            # Use any layout with a title
            layout = self.text_only
            fallback_config['use_body_as_subtitle'] = True

        elif slide_type == 'text_image':
            # Fall back to text_only with manual image placement
            layout = self.text_only
            if layout and layout.body_idx is not None:
                fallback_config['manual_image_placement'] = True
                fallback_config['resize_body_width_ratio'] = 0.48

        elif slide_type == 'text_image_top':
            # First try text_image (side-by-side), then text_only
            if self.text_image:
                layout = self.text_image
                fallback_config['use_side_by_side_for_stacked'] = True
            else:
                layout = self.text_only
                if layout and layout.body_idx is not None:
                    fallback_config['manual_image_placement'] = True
                    fallback_config['split_vertically'] = True

        return layout, fallback_config


def analyze_placeholder(ph) -> PlaceholderInfo:
    """
    Extract position and type information from a placeholder.

    Args:
        ph: A python-pptx placeholder shape

    Returns:
        PlaceholderInfo with position and type data
    """
    # Convert EMU to inches (914400 EMU per inch)
    EMU_PER_INCH = 914400

    left = ph.left / EMU_PER_INCH
    top = ph.top / EMU_PER_INCH
    width = ph.width / EMU_PER_INCH
    height = ph.height / EMU_PER_INCH

    ph_type = ph.placeholder_format.type

    # Get type name safely
    try:
        type_name = PP_PLACEHOLDER(ph_type).name if ph_type else "NONE"
    except ValueError:
        type_name = f"UNKNOWN_{ph_type}"

    return PlaceholderInfo(
        idx=ph.placeholder_format.idx,
        type=ph_type,
        type_name=type_name,
        left=left,
        top=top,
        width=width,
        height=height,
        center_x=left + width / 2,
        center_y=top + height / 2
    )


def detect_content_arrangement(
    placeholders: List[PlaceholderInfo],
    slide_width: float
) -> Tuple[str, Optional[int], Optional[int]]:
    """
    Detect whether content placeholders are arranged side-by-side or stacked.

    Args:
        placeholders: List of content placeholders (BODY or OBJECT type)
        slide_width: Width of the slide in inches

    Returns:
        Tuple of (arrangement, first_idx, second_idx) where:
        - arrangement is "side_by_side", "stacked", or "none"
        - first_idx is left/top placeholder index
        - second_idx is right/bottom placeholder index
    """
    # Filter to content placeholders only (BODY=2, OBJECT=7)
    content_phs = [p for p in placeholders if p.type in (2, 7)]

    if len(content_phs) < 2:
        return "none", None, None

    # Sort by position to identify arrangement
    # For side-by-side: sort by center_x
    # For stacked: sort by center_y

    # Check if side-by-side: Y centers within 0.5", X centers separated by > 1/3 slide width
    y_diff = abs(content_phs[0].center_y - content_phs[1].center_y)
    x_diff = abs(content_phs[0].center_x - content_phs[1].center_x)

    if y_diff < 0.5 and x_diff > slide_width / 3:
        # Side-by-side arrangement
        sorted_phs = sorted(content_phs, key=lambda p: p.center_x)
        return "side_by_side", sorted_phs[0].idx, sorted_phs[1].idx

    # Check if stacked: X centers within 1", Y centers separated by > 1"
    if x_diff < 1.0 and y_diff > 1.0:
        # Stacked arrangement
        sorted_phs = sorted(content_phs, key=lambda p: p.center_y)
        return "stacked", sorted_phs[0].idx, sorted_phs[1].idx

    # Default to side-by-side if ambiguous
    sorted_phs = sorted(content_phs, key=lambda p: p.center_x)
    return "side_by_side", sorted_phs[0].idx, sorted_phs[1].idx


def analyze_layout(layout, index: int, slide_width: float) -> LayoutCapabilities:
    """
    Analyze a single layout to determine its capabilities.

    Args:
        layout: A python-pptx slide layout
        index: The index of this layout in the presentation
        slide_width: Width of the slide in inches

    Returns:
        LayoutCapabilities describing what this layout can do
    """
    caps = LayoutCapabilities(index=index, name=layout.name)

    # Collect all placeholders
    all_phs: List[PlaceholderInfo] = []
    for ph in layout.placeholders:
        ph_info = analyze_placeholder(ph)
        all_phs.append(ph_info)
        caps._placeholders[ph_info.idx] = ph_info

    # Identify placeholder roles by type
    title_phs = [p for p in all_phs if p.type == 1]  # TITLE
    body_phs = [p for p in all_phs if p.type == 2]   # BODY
    object_phs = [p for p in all_phs if p.type == 7] # OBJECT

    # Title detection
    if title_phs:
        caps.has_title = True
        caps.title_idx = title_phs[0].idx

    # Analyze content structure
    content_phs = body_phs + object_phs

    if len(content_phs) == 1:
        # Single content area - distinguish between BODY and OBJECT types
        content_ph = content_phs[0]
        caps.body_idx = content_ph.idx  # Use body_idx for any single content placeholder

        if content_ph.type == 2:  # BODY type
            caps.has_single_body = True
        elif content_ph.type == 7:  # OBJECT type
            caps.has_single_object = True

        # Check if this is a subtitle (positioned below title, relatively small)
        if caps.has_title and title_phs:
            title_ph = title_phs[0]
            # Subtitle heuristic: below title, height < 2 inches
            if content_ph.top > title_ph.top + title_ph.height and content_ph.height < 2.5:
                caps.has_subtitle = True
                caps.subtitle_idx = content_ph.idx
                caps.has_single_body = False
                caps.has_single_object = False  # It's a subtitle, not content

    elif len(content_phs) >= 2:
        # Two or more content areas
        caps.has_two_content = True

        # Check if we have typed placeholders (BODY + OBJECT vs two OBJECTs)
        if body_phs and object_phs:
            caps.has_typed_placeholders = True

        # Detect arrangement
        arrangement, first_idx, second_idx = detect_content_arrangement(content_phs, slide_width)
        caps.content_arrangement = arrangement

        if arrangement == "side_by_side":
            caps.left_content_idx = first_idx
            caps.right_content_idx = second_idx
            # For typed layouts, identify which is body (text) vs object (image)
            if caps.has_typed_placeholders:
                for ph in content_phs:
                    if ph.type == 2:  # BODY
                        caps.body_idx = ph.idx
                    # Note: we keep left/right assignment based on position

        elif arrangement == "stacked":
            caps.top_content_idx = first_idx
            caps.bottom_content_idx = second_idx
            # For typed layouts, identify which is body (text) vs object (image)
            if caps.has_typed_placeholders:
                for ph in content_phs:
                    if ph.type == 2:  # BODY
                        caps.body_idx = ph.idx

    return caps


def discover_layouts(prs: Presentation, verbose: bool = True) -> LayoutMapping:
    """
    Discover and analyze all layouts in a presentation template.

    Args:
        prs: A python-pptx Presentation object
        verbose: Whether to print discovered layout information

    Returns:
        LayoutMapping with layouts assigned to slide types
    """
    # Get slide dimensions
    EMU_PER_INCH = 914400
    slide_width = prs.slide_width / EMU_PER_INCH
    slide_height = prs.slide_height / EMU_PER_INCH

    mapping = LayoutMapping()

    # Analyze all layouts
    for idx, layout in enumerate(prs.slide_layouts):
        caps = analyze_layout(layout, idx, slide_width)
        mapping.all_layouts.append(caps)

    # Assign layouts to slide types based on capabilities
    # Priority: prefer typed placeholders (BODY + OBJECT) over two OBJECTs

    for caps in mapping.all_layouts:
        # Title slide: has title + subtitle (or small body below title)
        if caps.has_title and (caps.has_subtitle or
            (caps.has_single_body and not mapping.title)):
            if mapping.title is None:
                mapping.title = caps

        # Section header: same as title for now
        if caps.has_title and caps.has_subtitle and mapping.section is None:
            mapping.section = caps

        # Text-only: has title + single full-width BODY
        if (caps.has_title and caps.has_single_body and
            not caps.has_subtitle and mapping.text_only is None):
            mapping.text_only = caps

        # Content-only: has title + single full-width OBJECT (for tables, images, etc.)
        if (caps.has_title and caps.has_single_object and
            not caps.has_subtitle and mapping.content_only is None):
            mapping.content_only = caps

        # Text+Image side-by-side: prefer typed (BODY left + OBJECT right)
        if caps.has_title and caps.has_two_content and caps.content_arrangement == "side_by_side":
            if caps.has_typed_placeholders:
                # Strongly prefer typed layout
                if mapping.text_image is None or not mapping.text_image.has_typed_placeholders:
                    mapping.text_image = caps
            elif mapping.text_image is None:
                mapping.text_image = caps

        # Text+Image stacked: prefer typed (OBJECT top + BODY bottom)
        if caps.has_title and caps.has_two_content and caps.content_arrangement == "stacked":
            if caps.has_typed_placeholders:
                # Strongly prefer typed layout
                if mapping.text_image_top is None or not mapping.text_image_top.has_typed_placeholders:
                    mapping.text_image_top = caps
            elif mapping.text_image_top is None:
                mapping.text_image_top = caps

    # text_diagram and text_table default to text_image and text_only respectively
    # (handled in get_layout_for_type)

    if verbose:
        print("Discovered layouts in template:")
        if mapping.title:
            print(f"  Title: {mapping.title.name} (index {mapping.title.index})")
        else:
            print("  Title: None (will use text_only fallback)")

        if mapping.text_only:
            print(f"  Text-only: {mapping.text_only.name} (index {mapping.text_only.index})")
        else:
            print("  Text-only: None")

        if mapping.content_only:
            print(f"  Content-only: {mapping.content_only.name} (index {mapping.content_only.index})")
            print(f"    Content placeholder idx={mapping.content_only.body_idx}")
        else:
            print("  Content-only: None (will use text_only fallback)")

        if mapping.text_image:
            print(f"  Text+Image (side-by-side): {mapping.text_image.name} (index {mapping.text_image.index})")
            if mapping.text_image.has_typed_placeholders:
                print(f"    Text placeholder idx={mapping.text_image.body_idx}, Image area idx={mapping.text_image.right_content_idx}")
        else:
            print("  Text+Image (side-by-side): None (will use text_only with manual placement)")

        if mapping.text_image_top:
            print(f"  Text+Image (stacked): {mapping.text_image_top.name} (index {mapping.text_image_top.index})")
            if mapping.text_image_top.has_typed_placeholders:
                print(f"    Image area idx={mapping.text_image_top.top_content_idx}, Text placeholder idx={mapping.text_image_top.body_idx}")
        else:
            print("  Text+Image (stacked): None (will use text_image or text_only fallback)")
        print()

    return mapping


def get_fallback_positioning(
    caps: LayoutCapabilities,
    slide_type: str,
    slide_width: float,
    slide_height: float
) -> Dict:
    """
    Calculate manual image positioning when ideal layout is unavailable.

    Args:
        caps: The layout capabilities being used (may be a fallback)
        slide_type: The intended slide type
        slide_width: Slide width in inches
        slide_height: Slide height in inches

    Returns:
        Dict with positioning configuration
    """
    config = {}

    if not caps or caps.body_idx is None:
        return config

    body_info = caps.get_placeholder_info(caps.body_idx)
    if not body_info:
        return config

    if slide_type == 'text_image' and caps.has_single_body:
        # Split body area horizontally for side-by-side
        config['resize_body_width'] = body_info.width * 0.48
        config['manual_image'] = {
            'left': body_info.left + body_info.width * 0.52,
            'top': body_info.top,
            'max_width': body_info.width * 0.48,
            'max_height': body_info.height
        }

    elif slide_type == 'text_image_top' and caps.has_single_body:
        # Split body area vertically for stacked
        config['manual_image'] = {
            'left': body_info.left,
            'top': body_info.top,
            'max_width': body_info.width,
            'max_height': body_info.height * 0.55
        }
        config['resize_body_top'] = body_info.top + body_info.height * 0.60
        config['resize_body_height'] = body_info.height * 0.40

    return config


# For testing/debugging
if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2:
        print("Usage: python layout_discovery.py <template.pptx>")
        sys.exit(1)

    prs = Presentation(sys.argv[1])
    mapping = discover_layouts(prs, verbose=True)

    print("\nDetailed layout analysis:")
    for caps in mapping.all_layouts:
        print(f"\nLayout {caps.index}: {caps.name}")
        print(f"  has_title={caps.has_title}, has_subtitle={caps.has_subtitle}")
        print(f"  has_single_body={caps.has_single_body}, has_two_content={caps.has_two_content}")
        print(f"  content_arrangement={caps.content_arrangement}")
        print(f"  has_typed_placeholders={caps.has_typed_placeholders}")
        if caps.title_idx is not None:
            print(f"  title_idx={caps.title_idx}")
        if caps.body_idx is not None:
            print(f"  body_idx={caps.body_idx}")
        if caps.left_content_idx is not None:
            print(f"  left_content_idx={caps.left_content_idx}, right_content_idx={caps.right_content_idx}")
        if caps.top_content_idx is not None:
            print(f"  top_content_idx={caps.top_content_idx}, bottom_content_idx={caps.bottom_content_idx}")
