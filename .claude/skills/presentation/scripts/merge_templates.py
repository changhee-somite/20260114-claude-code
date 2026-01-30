#!/usr/bin/env python3
"""
Merge layouts from skill template into official template.

This creates a hybrid template that combines:
- Official template styling and branding
- Rich layout options from the skill template (two-content layouts)
"""

import argparse
import copy
from pathlib import Path
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsmap, qn
import zipfile
import shutil
import tempfile
import os


def get_layout_info(prs):
    """Get information about all layouts in a presentation."""
    layouts = {}
    for i, layout in enumerate(prs.slide_master.slide_layouts):
        placeholders = []
        for ph in layout.placeholders:
            placeholders.append({
                'type': str(ph.placeholder_format.type),
                'idx': ph.placeholder_format.idx
            })
        layouts[layout.name] = {
            'index': i,
            'placeholders': placeholders
        }
    return layouts


def copy_layout_xml(source_pptx, target_pptx, layout_name, new_name=None):
    """
    Copy a layout from source to target presentation using XML manipulation.

    This is complex because layouts reference the slide master and theme.
    For simplicity, we'll recreate the layouts programmatically.
    """
    # This approach is complex - instead we'll document what needs manual work
    pass


def create_two_content_layout(prs, name, arrangement='horizontal'):
    """
    Create a two-content layout programmatically.

    Args:
        prs: Presentation object
        name: Layout name
        arrangement: 'horizontal' (side-by-side) or 'vertical' (stacked)
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.enum.shapes import PP_PLACEHOLDER

    # Get slide master
    master = prs.slide_master

    # We can't easily add new layouts programmatically in python-pptx
    # The library doesn't support creating new slide layouts
    # We need to use XML manipulation or manual work in PowerPoint

    return None


def analyze_and_report(official_path, skill_path):
    """Analyze both templates and report what needs to be merged."""
    official = Presentation(official_path)
    skill = Presentation(skill_path)

    print("=" * 60)
    print("TEMPLATE MERGE ANALYSIS")
    print("=" * 60)

    official_layouts = get_layout_info(official)
    skill_layouts = get_layout_info(skill)

    print("\n📋 Official Template Layouts:")
    for name, info in official_layouts.items():
        print(f"  [{info['index']}] {name}")

    print("\n📋 Skill Template Layouts:")
    for name, info in skill_layouts.items():
        print(f"  [{info['index']}] {name}")

    # Find layouts in skill that aren't in official
    missing = []
    for name in skill_layouts:
        if name not in official_layouts:
            missing.append(name)

    print("\n🔧 Layouts to add to official template:")
    for name in missing:
        info = skill_layouts[name]
        print(f"  - {name}")
        for ph in info['placeholders']:
            print(f"      Placeholder: {ph['type']} (idx={ph['idx']})")

    # Layouts in official but not in skill (new features)
    new_in_official = []
    for name in official_layouts:
        if name not in skill_layouts:
            new_in_official.append(name)

    if new_in_official:
        print("\n✨ New layouts from official template:")
        for name in new_in_official:
            print(f"  - {name}")

    return missing, new_in_official


def merge_via_xml(official_path, skill_path, output_path):
    """
    Merge templates by extracting and combining XML.

    PowerPoint files are ZIP archives containing XML files.
    We can extract, modify, and repackage them.
    """
    import re

    # Create temp directories
    with tempfile.TemporaryDirectory() as tmpdir:
        official_dir = Path(tmpdir) / "official"
        skill_dir = Path(tmpdir) / "skill"
        output_dir = Path(tmpdir) / "output"

        # Extract both
        with zipfile.ZipFile(official_path, 'r') as z:
            z.extractall(official_dir)
        with zipfile.ZipFile(skill_path, 'r') as z:
            z.extractall(skill_dir)

        # Copy official as base
        shutil.copytree(official_dir, output_dir)

        # Find the highest slideLayout number in official
        layouts_dir = output_dir / "ppt" / "slideLayouts"
        existing_layouts = list(layouts_dir.glob("slideLayout*.xml"))
        max_num = max(int(re.search(r'slideLayout(\d+)', f.name).group(1))
                      for f in existing_layouts)

        print(f"\n📁 Official template has {len(existing_layouts)} layouts (max: slideLayout{max_num}.xml)")

        # Identify layouts to copy from skill template
        skill_layouts_dir = skill_dir / "ppt" / "slideLayouts"

        # Parse skill presentation.xml to get layout names
        skill_pres = Presentation(skill_path)
        layouts_to_copy = [
            "Title & Content",
            "Two Content (horizontal)",
            "Text and Content (horizontal)",
            "Two Content (vertical)",
            "Content and Text (vertical)"
        ]

        # Find which slideLayoutN.xml corresponds to each layout name
        skill_layout_files = {}
        for i, layout in enumerate(skill_pres.slide_master.slide_layouts, start=1):
            if layout.name in layouts_to_copy:
                skill_layout_files[layout.name] = f"slideLayout{i}.xml"

        print(f"\n📋 Layouts to copy from skill template:")
        for name, filename in skill_layout_files.items():
            print(f"  - {name} ({filename})")

        # Copy layout files and update relationships
        # This is complex because we need to:
        # 1. Copy the slideLayoutN.xml files with new numbers
        # 2. Update [Content_Types].xml
        # 3. Update ppt/_rels/presentation.xml.rels
        # 4. Update ppt/slideLayouts/_rels/slideLayoutN.xml.rels
        # 5. Update ppt/slideMasters/_rels/slideMaster1.xml.rels
        # 6. Update ppt/slideMasters/slideMaster1.xml

        copied_layouts = {}
        new_num = max_num

        for name, src_filename in skill_layout_files.items():
            new_num += 1
            dst_filename = f"slideLayout{new_num}.xml"

            src_path = skill_layouts_dir / src_filename
            dst_path = layouts_dir / dst_filename

            if src_path.exists():
                # Read and copy the layout XML
                shutil.copy(src_path, dst_path)

                # Copy the .rels file too
                src_rels = skill_layouts_dir / "_rels" / f"{src_filename}.rels"
                dst_rels = layouts_dir / "_rels" / f"{dst_filename}.rels"
                if src_rels.exists():
                    # Read rels and update slideMaster reference
                    rels_content = src_rels.read_text()
                    # The relationship to slideMaster should point to ../slideMasters/slideMaster1.xml
                    # which is usually already correct
                    dst_rels.write_text(rels_content)

                copied_layouts[name] = {
                    'src': src_filename,
                    'dst': dst_filename,
                    'num': new_num
                }
                print(f"  ✓ Copied {src_filename} → {dst_filename}")

        # Update [Content_Types].xml
        content_types_path = output_dir / "[Content_Types].xml"
        ct_tree = etree.parse(str(content_types_path))
        ct_root = ct_tree.getroot()

        # Add Override entries for new layouts
        ns = {'ct': 'http://schemas.openxmlformats.org/package/2006/content-types'}
        for name, info in copied_layouts.items():
            override = etree.SubElement(ct_root, '{http://schemas.openxmlformats.org/package/2006/content-types}Override')
            override.set('PartName', f"/ppt/slideLayouts/{info['dst']}")
            override.set('ContentType', 'application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml')

        ct_tree.write(str(content_types_path), xml_declaration=True, encoding='UTF-8', standalone=True)
        print(f"\n✓ Updated [Content_Types].xml")

        # Update slideMaster1.xml.rels to include new layouts
        master_rels_path = output_dir / "ppt" / "slideMasters" / "_rels" / "slideMaster1.xml.rels"
        master_rels_tree = etree.parse(str(master_rels_path))
        master_rels_root = master_rels_tree.getroot()

        # Find highest rId
        max_rid = 0
        for rel in master_rels_root:
            rid = rel.get('Id', '')
            if rid.startswith('rId'):
                try:
                    num = int(rid[3:])
                    max_rid = max(max_rid, num)
                except ValueError:
                    pass

        # Add relationships for new layouts
        rel_ns = 'http://schemas.openxmlformats.org/package/2006/relationships'
        layout_type = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout'

        for name, info in copied_layouts.items():
            max_rid += 1
            rel = etree.SubElement(master_rels_root, '{%s}Relationship' % rel_ns)
            rel.set('Id', f'rId{max_rid}')
            rel.set('Type', layout_type)
            rel.set('Target', f'../slideLayouts/{info["dst"]}')
            info['rId'] = f'rId{max_rid}'

        master_rels_tree.write(str(master_rels_path), xml_declaration=True, encoding='UTF-8', standalone=True)
        print(f"✓ Updated slideMaster1.xml.rels")

        # Update slideMaster1.xml to reference new layouts
        master_path = output_dir / "ppt" / "slideMasters" / "slideMaster1.xml"
        master_tree = etree.parse(str(master_path))
        master_root = master_tree.getroot()

        # Find sldLayoutIdLst element
        p_ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        r_ns = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

        sld_layout_id_lst = master_root.find('.//{%s}sldLayoutIdLst' % p_ns)
        if sld_layout_id_lst is not None:
            # Find highest id
            max_id = 0
            for sld_layout_id in sld_layout_id_lst:
                id_val = sld_layout_id.get('id', '0')
                try:
                    max_id = max(max_id, int(id_val))
                except ValueError:
                    pass

            # Add new layout references
            for name, info in copied_layouts.items():
                max_id += 1
                sld_layout_id = etree.SubElement(sld_layout_id_lst, '{%s}sldLayoutId' % p_ns)
                sld_layout_id.set('id', str(max_id))
                sld_layout_id.set('{%s}id' % r_ns, info['rId'])

            master_tree.write(str(master_path), xml_declaration=True, encoding='UTF-8', standalone=True)
            print(f"✓ Updated slideMaster1.xml")

        # Repackage as PPTX
        output_pptx = Path(output_path)
        if output_pptx.exists():
            output_pptx.unlink()

        with zipfile.ZipFile(output_pptx, 'w', zipfile.ZIP_DEFLATED) as z:
            for root, dirs, files in os.walk(output_dir):
                for file in files:
                    file_path = Path(root) / file
                    arcname = file_path.relative_to(output_dir)
                    z.write(file_path, arcname)

        print(f"\n✅ Created merged template: {output_path}")

        return copied_layouts


def main():
    parser = argparse.ArgumentParser(description='Merge PowerPoint templates')
    parser.add_argument('--official', required=True, help='Path to official template')
    parser.add_argument('--skill', required=True, help='Path to skill template')
    parser.add_argument('--output', required=True, help='Output path for merged template')
    parser.add_argument('--analyze-only', action='store_true', help='Only analyze, do not merge')

    args = parser.parse_args()

    if args.analyze_only:
        analyze_and_report(args.official, args.skill)
    else:
        analyze_and_report(args.official, args.skill)
        merge_via_xml(args.official, args.skill, args.output)


if __name__ == '__main__':
    main()
