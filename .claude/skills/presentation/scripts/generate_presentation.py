#!/usr/bin/env python3
"""
Unified Presentation Generator

Generates PowerPoint presentations from PRESENTATION.md with integrated
text, image, and diagram handling.

Features:
- Parses PRESENTATION.md to extract slide content
- Detects slides with figures, diagrams, or text-only
- Assigns appropriate layouts (split for images/diagrams, full-width for text)
- Inserts images with proper positioning
- Renders vector diagrams using python-pptx shapes
- Resizes text placeholders to avoid overlap

Usage:
    python generate_presentation.py \\
        --template template.pptx \\
        --source docs/PRESENTATION.md \\
        --figures figures/ \\
        --output final.pptx

Author: Claude Code
Date: 2026-01-17
"""

import argparse
import json
import os
import re
import sys
from dataclasses import dataclass, field
from pathlib import Path
from typing import List, Optional, Dict, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree
from PIL import Image

# Import diagram renderer
from diagram_renderer import render_diagram, parse_diagram_block, DiagramConfig

# Import table renderer
from table_renderer import render_table, parse_table_block, parse_markdown_table, TableConfig, TableStyle

# Import layout discovery
from layout_discovery import (
    discover_layouts, LayoutMapping, LayoutCapabilities,
    get_fallback_positioning
)

# Default template path (relative to this script's directory)
DEFAULT_TEMPLATE = Path(__file__).parent.parent / "templates" / "cellularintelligence.pptx"

# Template font sizes (from slideMaster1.xml bodyStyle)
# These are the default font sizes for each bullet level in the template
TEMPLATE_FONT_SIZES = {
    0: 24,  # lvl1pPr - main bullet
    1: 18,  # lvl2pPr - sub-bullet
    2: 15,  # lvl3pPr
    3: 12,  # lvl4pPr
    4: 10,  # lvl5pPr
}

# Template paragraph spacing (spcBef - space before each paragraph, in points)
# From slideMaster1.xml bodyStyle
TEMPLATE_PARAGRAPH_SPACING = {
    0: 12,   # lvl1pPr
    1: 12,   # lvl2pPr
    2: 12,   # lvl3pPr
    3: 29.5, # lvl4pPr
    4: 29.5, # lvl5pPr
}

MIN_FONT_SIZE = 10  # Don't shrink fonts below this size


def disable_bullet(paragraph):
    """
    Disable the bullet for a paragraph by adding <a:buNone/> to the XML.

    python-pptx doesn't have a built-in way to disable bullets, so we
    manipulate the underlying XML directly.
    """
    # Get or create paragraph properties
    pPr = paragraph._p.get_or_add_pPr()

    # Remove any existing bullet elements
    for child in list(pPr):
        tag = etree.QName(child.tag).localname
        if tag.startswith('bu'):
            pPr.remove(child)

    # Add buNone to disable bullet
    etree.SubElement(pPr, qn('a:buNone'))


@dataclass
class ContentItem:
    """Represents a single content item with formatting metadata."""
    text: str
    is_code: bool = False  # True if this is a code block line
    code_language: str = ""  # Language hint (e.g., "python", "bash")


@dataclass
class SlideContent:
    """Represents parsed content for a single slide."""
    number: int
    title: str = ""
    content: List[str] = field(default_factory=list)  # Plain strings for backward compat
    content_items: List[ContentItem] = field(default_factory=list)  # Rich content items
    figure_path: Optional[str] = None
    figure_position: str = "right"  # right (default) or top
    diagram_spec: Optional[Dict] = None  # Parsed diagram specification
    table_spec: Optional[Dict] = None    # Parsed table specification
    notes: Optional[str] = None
    slide_type: str = "text_only"  # text_only, text_image, text_image_top, text_diagram, text_table, section, title


@dataclass
class LayoutConfig:
    """Configuration for slide layouts."""
    text_only_idx: int = 1
    text_image_idx: int = 1  # Same as text_only if no picture placeholder
    section_idx: int = 0
    title_idx: int = 0

    # Text placeholder dimensions for image slides (in inches)
    text_width: float = 6.5  # Width of text area when image present
    image_left: float = 7.0   # Left position of image
    image_width: float = 5.5  # Width of image
    image_top: float = 1.5    # Top position of image (below title)

    # Image-on-top layout dimensions (in inches)
    image_top_layout_image_top: float = 1.5     # Image starts below title
    image_top_layout_image_height: float = 4.0  # Max image height
    image_top_layout_text_top: float = 5.7      # Text starts below image
    image_top_layout_text_height: float = 1.5   # Reduced text area height


def parse_presentation_md(md_path: Path) -> List[SlideContent]:
    """Parse PRESENTATION.md and extract slide content."""

    with open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()

    slides = []

    # Split by slide headers (### Slide N: or ### Slide AN:)
    slide_pattern = r'###\s+Slide\s+(\d+|A\d+):\s*(.+?)(?=###\s+Slide|\Z)'
    matches = re.findall(slide_pattern, content, re.DOTALL)

    for slide_num_str, slide_content in matches:
        slide = SlideContent(
            number=int(slide_num_str) if slide_num_str.isdigit() else 100 + int(slide_num_str[1:])
        )

        # Extract title
        title_match = re.search(r'\*\*Title\*\*:\s*(.+?)(?:\n|$)', slide_content)
        if title_match:
            slide.title = title_match.group(1).strip()

        # Extract type if specified
        type_match = re.search(r'\*\*Type\*\*:\s*(.+?)(?:\n|$)', slide_content)
        if type_match:
            type_str = type_match.group(1).strip().lower()
            if 'section' in type_str:
                slide.slide_type = 'section'

        # Extract content (bullets from **Content**: section)
        content_match = re.search(
            r'\*\*Content\*\*:\s*\n((?:(?!\*\*(?:Figure|Notes|Source)\*\*).)+)',
            slide_content,
            re.DOTALL
        )
        if content_match:
            content_text = content_match.group(1)
            # Parse bullets and other content, PRESERVING indentation for hierarchy
            # Also detect and extract markdown tables and code blocks from content
            bullets = []
            content_items = []
            table_lines = []  # Collect consecutive table rows
            in_code_block = False
            code_language = ""

            for line in content_text.split('\n'):
                line_stripped = line.strip()

                # Handle code fence markers
                if line_stripped.startswith('```'):
                    if not in_code_block:
                        # Starting a code block
                        in_code_block = True
                        code_language = line_stripped[3:].strip()  # e.g., "python", "bash"
                    else:
                        # Ending a code block
                        in_code_block = False
                        code_language = ""
                    continue

                # If inside code block, preserve the line as code
                if in_code_block:
                    # Preserve original indentation for code
                    bullets.append(line.rstrip())  # Keep leading spaces, strip trailing
                    content_items.append(ContentItem(
                        text=line.rstrip(),
                        is_code=True,
                        code_language=code_language
                    ))
                    continue

                if not line_stripped:
                    continue

                # Check for markdown table rows
                if line_stripped.startswith('|') and line_stripped.endswith('|'):
                    table_lines.append(line_stripped)
                    continue
                else:
                    # If we were collecting table lines and hit non-table, finalize the table
                    if table_lines and not slide.table_spec:
                        # Parse the collected table
                        table_md = '\n'.join(table_lines)
                        slide.table_spec = parse_markdown_table(table_md)
                        if slide.table_spec:
                            slide.slide_type = 'text_table'
                        table_lines = []

                # Check if it's a bullet line (but not a **bold** line)
                if line_stripped.startswith(('-', '*', '•')) and not line_stripped.startswith('**'):
                    # Count leading spaces BEFORE stripping to preserve hierarchy
                    leading_spaces = len(line) - len(line.lstrip())
                    # Remove the bullet marker from the stripped version
                    bullet_text = re.sub(r'^[-*•]\s*', '', line_stripped)
                    if bullet_text:
                        # Preserve indentation by prepending spaces (2 spaces per level)
                        text = ' ' * leading_spaces + bullet_text
                        bullets.append(text)
                        content_items.append(ContentItem(text=text, is_code=False))
                elif line_stripped.startswith('>'):
                    # Quote
                    quote_text = line_stripped.lstrip('> ')
                    text = f'"{quote_text}"'
                    bullets.append(text)
                    content_items.append(ContentItem(text=text, is_code=False))
                elif not line_stripped.startswith('**'):
                    # Non-bullet content - preserve indentation
                    leading_spaces = len(line) - len(line.lstrip())
                    text = ' ' * leading_spaces + line_stripped
                    bullets.append(text)
                    content_items.append(ContentItem(text=text, is_code=False))

            # Handle table at end of content
            if table_lines and not slide.table_spec:
                table_md = '\n'.join(table_lines)
                slide.table_spec = parse_markdown_table(table_md)
                if slide.table_spec:
                    slide.slide_type = 'text_table'

            slide.content = bullets
            slide.content_items = content_items

        # Extract figure path
        figure_match = re.search(
            r'\*\*Figure\*\*:\s*\[.+?\]\(([^)]+)\)',
            slide_content
        )
        if figure_match:
            slide.figure_path = figure_match.group(1)
            slide.slide_type = 'text_image'

        # Extract figure position (top or right, default is right)
        figure_pos_match = re.search(
            r'\*\*Figure Position\*\*:\s*(\w+)',
            slide_content
        )
        if figure_pos_match:
            pos = figure_pos_match.group(1).lower()
            if pos == 'top':
                slide.figure_position = 'top'
                slide.slide_type = 'text_image_top'

        # Extract diagram specification
        diagram_match = re.search(
            r'\*\*Diagram\*\*:\s*\n```diagram\n(.+?)```',
            slide_content,
            re.DOTALL
        )
        if diagram_match:
            diagram_block = diagram_match.group(1)
            slide.diagram_spec = parse_diagram_block(diagram_block)
            slide.slide_type = 'text_diagram'

        # Extract table specification (markdown table or ```table block)
        table_match = re.search(
            r'\*\*Table\*\*:\s*\n```table\n(.+?)```',
            slide_content,
            re.DOTALL
        )
        if table_match:
            table_block = table_match.group(1)
            slide.table_spec = parse_table_block(table_block)
            slide.slide_type = 'text_table'
        else:
            # Check for inline markdown table in content
            md_table_match = re.search(
                r'\*\*Table\*\*:\s*\n(\|.+?\|(?:\n\|.+?\|)+)',
                slide_content,
                re.DOTALL
            )
            if md_table_match:
                md_table = md_table_match.group(1)
                slide.table_spec = parse_markdown_table(md_table)
                slide.slide_type = 'text_table'

        # Extract notes
        notes_match = re.search(r'\*\*Notes\*\*:\s*(.+?)(?:\n\*\*|\Z)', slide_content, re.DOTALL)
        if notes_match:
            slide.notes = notes_match.group(1).strip()

        # Determine slide type (only override if not already set to a specific type)
        if slide.number == 1:
            slide.slide_type = 'title'
        elif slide.slide_type not in ('text_image', 'text_image_top', 'text_diagram', 'text_table', 'section'):
            slide.slide_type = 'text_only'

        slides.append(slide)

    return slides


def emu_to_inches(emu) -> float:
    """Convert EMUs to inches."""
    return emu / 914400


def get_image_dimensions(image_path: Path) -> Tuple[int, int]:
    """Get image dimensions in pixels."""
    with Image.open(image_path) as img:
        return img.size


def find_placeholder(slide, placeholder_type=None, idx=None):
    """Find a placeholder by type or index."""
    for shape in slide.placeholders:
        if idx is not None and shape.placeholder_format.idx == idx:
            return shape
        if placeholder_type is not None:
            if shape.placeholder_format.type == placeholder_type:
                return shape
    return None


def find_body_placeholder(slide):
    """Find the body/content placeholder."""
    for shape in slide.placeholders:
        ph_type = shape.placeholder_format.type
        if ph_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
            return shape
    return None


def find_content_placeholder_by_index(slide, idx: int):
    """Find a content placeholder by its index."""
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            return shape
    return None


def hide_content_placeholders(slide, exclude_idx: int = None):
    """
    Hide all content placeholders on a slide by moving them off-screen.

    This is used when we place content (like tables) directly on the slide
    without using the placeholder system, to prevent empty placeholder boxes
    from showing.

    Args:
        slide: The slide object
        exclude_idx: Placeholder index to NOT hide (e.g., if we're using it for text)
    """
    for shape in slide.placeholders:
        ph_type = shape.placeholder_format.type
        ph_idx = shape.placeholder_format.idx

        # Skip title placeholder (idx 0) and any excluded placeholder
        if ph_idx == 0:
            continue
        if exclude_idx is not None and ph_idx == exclude_idx:
            continue

        # Hide content/body/object placeholders
        if ph_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
            shape.left = Inches(-10)  # Move off-slide
            shape.top = Inches(-10)
            shape.width = Inches(0.01)
            shape.height = Inches(0.01)


def get_placeholder_bounds(slide, placeholder_idx: int) -> Optional[Dict]:
    """
    Get the bounds of a placeholder by index.

    Returns:
        Dict with 'left', 'top', 'width', 'height' in inches, or None if not found
    """
    placeholder = find_content_placeholder_by_index(slide, placeholder_idx)
    if not placeholder:
        return None

    return {
        'left': emu_to_inches(placeholder.left),
        'top': emu_to_inches(placeholder.top),
        'width': emu_to_inches(placeholder.width),
        'height': emu_to_inches(placeholder.height),
    }


def insert_image_into_placeholder(slide, image_path: Path, placeholder_idx: int) -> Optional[object]:
    """
    Insert an image into a specific placeholder, fitting within its bounds.

    The placeholder is hidden after the image is inserted.

    Args:
        slide: The slide object
        image_path: Path to the image file
        placeholder_idx: Index of the OBJECT placeholder to use

    Returns:
        Picture shape or None if failed
    """
    placeholder = find_content_placeholder_by_index(slide, placeholder_idx)
    if not placeholder:
        print(f"    Warning: Placeholder [{placeholder_idx}] not found")
        return None

    if not image_path.exists():
        print(f"    Warning: Image not found: {image_path}")
        return None

    # Get placeholder bounds
    ph_left = emu_to_inches(placeholder.left)
    ph_top = emu_to_inches(placeholder.top)
    ph_width = emu_to_inches(placeholder.width)
    ph_height = emu_to_inches(placeholder.height)

    # Get image dimensions for aspect ratio
    img_width, img_height = get_image_dimensions(image_path)
    aspect_ratio = img_width / img_height

    # Calculate size to fit within placeholder while maintaining aspect ratio
    if ph_width / ph_height > aspect_ratio:
        # Height is the constraint
        height = ph_height
        width = height * aspect_ratio
    else:
        # Width is the constraint
        width = ph_width
        height = width / aspect_ratio

    # Center the image within the placeholder
    left = ph_left + (ph_width - width) / 2
    top = ph_top + (ph_height - height) / 2

    # Add the picture
    pic = slide.shapes.add_picture(
        str(image_path),
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height)
    )

    # Hide the placeholder
    placeholder.left = Inches(-10)
    placeholder.top = Inches(-10)
    placeholder.width = Inches(0.01)
    placeholder.height = Inches(0.01)

    return pic


def insert_table_into_placeholder(slide, table_spec: Dict, placeholder_idx: int,
                                   style: TableStyle = None) -> Optional[object]:
    """
    Insert a table into a specific placeholder, fitting within its bounds.

    The placeholder is hidden after the table is inserted.

    Args:
        slide: The slide object
        table_spec: Dict with 'headers' and 'rows'
        placeholder_idx: Index of the OBJECT placeholder to use
        style: Optional TableStyle for appearance

    Returns:
        Table shape or None if failed
    """
    placeholder = find_content_placeholder_by_index(slide, placeholder_idx)
    if not placeholder:
        print(f"    Warning: Placeholder [{placeholder_idx}] not found")
        return None

    # Get placeholder bounds
    ph_left = emu_to_inches(placeholder.left)
    ph_top = emu_to_inches(placeholder.top)
    ph_width = emu_to_inches(placeholder.width)
    ph_height = emu_to_inches(placeholder.height)

    # Create table config from placeholder bounds
    config = TableConfig(
        left=ph_left,
        top=ph_top,
        width=ph_width
    )

    # Render the table
    table_shape = render_table(slide, table_spec, config, style)

    # Hide the placeholder
    placeholder.left = Inches(-10)
    placeholder.top = Inches(-10)
    placeholder.width = Inches(0.01)
    placeholder.height = Inches(0.01)

    return table_shape


def fill_body_placeholder(slide, content: List[str], placeholder_idx: int,
                          font_scale: float = 1.0,
                          content_items: List[ContentItem] = None) -> Tuple[bool, float]:
    """
    Fill a BODY placeholder with bullet content.

    Args:
        slide: The slide object
        content: List of text items (backward compat)
        placeholder_idx: Index of the BODY placeholder (typically 1 or 13)
        font_scale: Font scale factor (1.0 = no scaling)
        content_items: List of ContentItem with formatting metadata (preferred)

    Returns:
        Tuple of (was_scaled, scale_factor)
    """
    placeholder = find_content_placeholder_by_index(slide, placeholder_idx)
    if not placeholder or not hasattr(placeholder, 'text_frame'):
        # Fall back to finding any BODY placeholder
        placeholder = find_body_placeholder(slide)
        if not placeholder:
            return False, 1.0

    # Get dimensions for overflow check
    width = emu_to_inches(placeholder.width)
    height = emu_to_inches(placeholder.height)

    # Check for overflow
    was_scaled = False
    if font_scale == 1.0:
        overflow, scale = estimate_content_overflow(content, width, height)
        if overflow:
            font_scale = scale
            was_scaled = True

    tf = placeholder.text_frame
    tf.clear()

    # Use content_items if available, otherwise fall back to content list
    items_to_render = content_items if content_items else [ContentItem(text=c) for c in content]

    for i, item in enumerate(items_to_render):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Get text and code flag from ContentItem
        item_text = item.text if isinstance(item, ContentItem) else item
        is_code = item.is_code if isinstance(item, ContentItem) else False

        # Skip figure/table references
        if item_text.startswith('[Figure:') or item_text.startswith('**Figure**'):
            continue

        # For code blocks, disable bullet and use monospace
        if is_code:
            level = 0
            disable_bullet(p)
            apply_formatted_text(p, item_text, is_code=True)
        else:
            # Detect indentation level and clean the text
            level, cleaned_text, is_numbered = detect_indent_level(item_text)
            p.level = level

            # For numbered items, disable the bullet
            if is_numbered:
                disable_bullet(p)

            # Apply formatted text (handles **bold**, [link](url), `code` patterns)
            apply_formatted_text(p, cleaned_text)

        # Apply scaled font size if needed
        if font_scale < 1.0:
            base_size = TEMPLATE_FONT_SIZES.get(level, 18)
            new_size = max(MIN_FONT_SIZE, int(base_size * font_scale))
            for run in p.runs:
                run.font.size = Pt(new_size)

    return was_scaled, font_scale


def insert_image_in_placeholder(slide, image_path: Path, placeholder_idx: int):
    """
    Insert an image into a placeholder's bounds, then hide the placeholder.

    This properly fills the designated area instead of adding a free-floating image.
    """
    placeholder = find_content_placeholder_by_index(slide, placeholder_idx)
    if not placeholder:
        return None

    # Get placeholder bounds
    left = placeholder.left
    top = placeholder.top
    ph_width = placeholder.width
    ph_height = placeholder.height

    # Get image dimensions for aspect ratio
    img_width, img_height = get_image_dimensions(image_path)
    aspect_ratio = img_width / img_height

    # Calculate size to fit within placeholder while maintaining aspect ratio
    ph_width_inches = ph_width / 914400
    ph_height_inches = ph_height / 914400

    if ph_width_inches / ph_height_inches > aspect_ratio:
        # Height is the constraint
        height = ph_height_inches
        width = height * aspect_ratio
    else:
        # Width is the constraint
        width = ph_width_inches
        height = width / aspect_ratio

    # Center the image within the placeholder
    left_inches = (left / 914400) + (ph_width_inches - width) / 2
    top_inches = (top / 914400) + (ph_height_inches - height) / 2

    # Add the picture
    pic = slide.shapes.add_picture(
        str(image_path),
        Inches(left_inches),
        Inches(top_inches),
        Inches(width),
        Inches(height)
    )

    # Hide the placeholder by making it tiny and moving it off-slide
    # We can't delete placeholders easily, but we can make them invisible
    placeholder.left = Inches(-10)  # Move off-slide
    placeholder.top = Inches(-10)
    placeholder.width = Inches(0.01)
    placeholder.height = Inches(0.01)

    return pic


def fill_content_placeholder(slide, content: List[str], placeholder_idx: int,
                             font_scale: float = 1.0,
                             content_items: List[ContentItem] = None) -> Tuple[bool, float]:
    """
    Fill a specific content placeholder by index with bullet content.

    Args:
        slide: The slide object
        content: List of text items for the placeholder (backward compat)
        placeholder_idx: Index of the placeholder to fill
        font_scale: Font scale factor (1.0 = no scaling, < 1.0 = shrink)
        content_items: List of ContentItem with formatting metadata (preferred)

    Returns:
        Tuple of (was_scaled: bool, scale_factor: float)
    """
    placeholder = find_content_placeholder_by_index(slide, placeholder_idx)
    if not placeholder or not hasattr(placeholder, 'text_frame'):
        return False, 1.0

    # Get placeholder dimensions for overflow check
    placeholder_width = emu_to_inches(placeholder.width)
    placeholder_height = emu_to_inches(placeholder.height)

    # Check for overflow if no scale factor provided
    was_scaled = False
    if font_scale == 1.0:
        overflow, scale = estimate_content_overflow(content, placeholder_width, placeholder_height)
        if overflow:
            font_scale = scale
            was_scaled = True

    tf = placeholder.text_frame
    tf.clear()

    # Use content_items if available, otherwise fall back to content list
    items_to_render = content_items if content_items else [ContentItem(text=c) for c in content]

    for i, item in enumerate(items_to_render):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Get text and code flag from ContentItem
        item_text = item.text if isinstance(item, ContentItem) else item
        is_code = item.is_code if isinstance(item, ContentItem) else False

        # Skip figure references
        if item_text.startswith('[Figure:') or item_text.startswith('**Figure**'):
            continue

        # For code blocks, disable bullet and use monospace
        if is_code:
            level = 0
            disable_bullet(p)
            apply_formatted_text(p, item_text, is_code=True)
        else:
            # Detect indentation level and clean the text
            level, cleaned_text, is_numbered = detect_indent_level(item_text)
            p.level = level

            # For numbered items, disable the bullet
            if is_numbered:
                disable_bullet(p)

            # Apply formatted text (handles **bold**, [link](url), `code` patterns)
            apply_formatted_text(p, cleaned_text)

        # Apply scaled font size if needed
        if font_scale < 1.0:
            base_size = TEMPLATE_FONT_SIZES.get(level, 18)
            new_size = max(MIN_FONT_SIZE, int(base_size * font_scale))
            for run in p.runs:
                run.font.size = Pt(new_size)

    return was_scaled, font_scale


def apply_formatted_text(paragraph, text: str, is_code: bool = False, monospace_font: str = "Consolas"):
    """
    Apply text to a paragraph, parsing markdown patterns:
    - **bold** for bold text
    - [text](url) for hyperlinks
    - `code` for inline code (monospace)

    Creates separate runs for each formatting segment.

    Args:
        paragraph: python-pptx paragraph object
        text: Text that may contain markdown patterns
        is_code: If True, render entire text as monospace code
        monospace_font: Font name for code/monospace text
    """
    # Clear existing runs
    paragraph.clear()

    # If this is a code block line, apply monospace to entire text
    if is_code:
        run = paragraph.add_run()
        run.text = text
        run.font.name = monospace_font
        return

    # Combined pattern to find all formatting in order:
    # - **bold** text
    # - [text](url) markdown links
    # - `inline code`
    combined_pattern = re.compile(
        r'(\*\*(.+?)\*\*)|'          # Group 1-2: bold
        r'(\[([^\]]+)\]\(([^)]+)\))|'  # Group 3-5: markdown link [text](url)
        r'(`([^`]+)`)'                 # Group 6-7: inline code
    )

    last_end = 0
    for match in combined_pattern.finditer(text):
        # Add plain text before this match
        if match.start() > last_end:
            run = paragraph.add_run()
            run.text = text[last_end:match.start()]

        if match.group(1):  # Bold text
            run = paragraph.add_run()
            run.text = match.group(2)
            run.font.bold = True
        elif match.group(3):  # Markdown link [text](url)
            link_text = match.group(4)
            link_url = match.group(5)
            run = paragraph.add_run()
            run.text = link_text
            # Add hyperlink
            run.hyperlink.address = link_url
        elif match.group(6):  # Inline code `text`
            run = paragraph.add_run()
            run.text = match.group(7)
            run.font.name = monospace_font

        last_end = match.end()

    # Add any remaining plain text
    if last_end < len(text):
        run = paragraph.add_run()
        run.text = text[last_end:]


def detect_indent_level(text: str) -> Tuple[int, str, bool]:
    """
    Detect indentation level from leading whitespace.

    Args:
        text: Text that may have leading spaces indicating hierarchy

    Returns:
        Tuple of (level, cleaned_text, is_numbered) where:
        - level is 0-3
        - cleaned_text is the text to display
        - is_numbered indicates if this is a numbered list item (no bullet needed)
    """
    # Count leading spaces
    stripped = text.lstrip()
    leading_spaces = len(text) - len(stripped)

    # 2-4 spaces = 1 level of indentation
    level = min(3, leading_spaces // 2)  # Cap at level 3

    # Check for numbered list items (e.g., "1. ", "2. ", "10. ")
    numbered_match = re.match(r'^(\d+)\.\s+(.+)$', stripped)
    if numbered_match:
        # Keep the number as part of the text, mark as numbered
        return level, stripped, True

    # Check for sub-bullets (e.g., "- item" or "* item")
    # BUT NOT **bold** markdown patterns
    if stripped.startswith(('-', '•')):
        # Remove the bullet marker (dash or bullet only, not asterisk)
        cleaned = re.sub(r'^[-•]\s*', '', stripped)
        return level, cleaned, False
    elif stripped.startswith('*') and not stripped.startswith('**'):
        # Single asterisk bullet, not bold marker
        cleaned = re.sub(r'^\*\s*', '', stripped)
        return level, cleaned, False

    return level, stripped, False


def estimate_content_overflow(content: List[str], width_inches: float, height_inches: float) -> Tuple[bool, float]:
    """
    Estimate if content will overflow the placeholder and calculate scale factor.

    This function proactively checks if the content will fit within the given
    dimensions using the template's font sizes and paragraph spacing, and returns
    a scale factor to apply if shrinking is needed.

    Args:
        content: List of text items (may include indentation for hierarchy)
        width_inches: Available width in inches
        height_inches: Available height in inches

    Returns:
        Tuple of (will_overflow: bool, scale_factor: float)
        scale_factor is 1.0 if no scaling needed, < 1.0 to shrink
    """
    if not content:
        return False, 1.0

    # Calculate available space with margins
    effective_height_pt = height_inches * 72 * 0.85  # 15% margin for top/bottom
    effective_width_pt = width_inches * 72 * 0.90   # 10% margin for left/right

    # Filter out figure references
    items = [item for item in content
             if not item.startswith('[Figure:') and not item.startswith('**Figure**')]

    if not items:
        return False, 1.0

    # Calculate space needed based on actual content levels and template spacing
    # The slide master uses: lnSpc=100% (line height = font size) + spcBef (space before)
    lines_needed_pt = 0

    for i, item in enumerate(items):
        level, text, _ = detect_indent_level(item)
        font_size = TEMPLATE_FONT_SIZES.get(level, 18)
        para_spacing = TEMPLATE_PARAGRAPH_SPACING.get(level, 12)

        # Line height = font size (since lnSpc=100%)
        line_height = font_size

        # Estimate text wrapping - how many lines this item needs
        char_width = font_size * 0.5  # average character width

        # Account for bullet indent reducing available width
        indent_reduction = level * 0.3 * 72  # ~0.3" per indent level
        adjusted_width = max(72, effective_width_pt - indent_reduction)  # min 1 inch
        chars_per_line = max(1, adjusted_width / char_width)

        text_lines = max(1, len(text) / chars_per_line)

        # Add space: font height for all lines + space before (except first paragraph)
        lines_needed_pt += text_lines * line_height
        if i > 0:
            lines_needed_pt += para_spacing

    # Check if content fits
    if lines_needed_pt <= effective_height_pt:
        return False, 1.0

    # Calculate scale factor needed to fit
    # Note: scaling fonts doesn't scale paragraph spacing, so we need to be more aggressive
    # Estimate how much we can save by scaling fonts only
    def calc_item_height(item):
        level, text, _ = detect_indent_level(item)
        font_size = TEMPLATE_FONT_SIZES.get(level, 18)
        indent_reduction = level * 0.3 * 72
        adjusted_width = max(72, effective_width_pt - indent_reduction)
        chars_per_line = max(1, adjusted_width / (font_size * 0.5))
        text_lines = max(1, len(text) / chars_per_line)
        return text_lines * font_size

    font_only_height = sum(calc_item_height(item) for item in items)
    spacing_height = lines_needed_pt - font_only_height

    # Target: effective_height = scaled_font_height + spacing_height
    # scaled_font_height = font_only_height * scale
    # scale = (effective_height - spacing_height) / font_only_height
    if font_only_height > 0:
        scale_factor = (effective_height_pt - spacing_height) / font_only_height * 0.95
    else:
        scale_factor = effective_height_pt / lines_needed_pt * 0.95

    # Don't shrink below 50% (would be illegible)
    scale_factor = max(0.5, min(1.0, scale_factor))

    return True, scale_factor


def fill_title(slide, title_text: str):
    """
    Fill the title placeholder with bold formatting.

    Uses the slide master's font sizing for consistency.
    """
    title_shape = slide.shapes.title
    if not title_shape:
        return

    title_shape.text = title_text

    # Apply bold formatting to all runs in the title
    # Font size is inherited from the slide master for consistency
    for paragraph in title_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.bold = True


def fill_subtitle(slide, subtitle_text: str):
    """
    Fill the subtitle/body placeholder on a title slide.
    """
    # On title slides, the body placeholder is used for subtitle
    body_shape = find_body_placeholder(slide)
    if body_shape:
        body_shape.text = subtitle_text


def fill_body(slide, content: List[str], resize_width: float = None,
              font_scale: float = 1.0,
              content_items: List[ContentItem] = None) -> Tuple[bool, float]:
    """
    Fill the body placeholder with bullet content, with optional font scaling.

    Args:
        slide: The slide object
        content: List of text items for the body (backward compat)
        resize_width: Optional width to resize the placeholder to
        font_scale: Font scale factor (1.0 = no scaling, < 1.0 = shrink)
        content_items: List of ContentItem with formatting metadata (preferred)

    Returns:
        Tuple of (was_scaled: bool, scale_factor: float)
    """
    body_shape = find_body_placeholder(slide)
    if not body_shape:
        return False, 1.0

    # IMPORTANT: Save ALL original position/size before any modifications
    # Setting only some properties creates incomplete XML (python-pptx bug)
    original_left = body_shape.left
    original_top = body_shape.top
    original_width = body_shape.width
    original_height = body_shape.height

    # Determine effective dimensions for overflow check
    effective_width = resize_width if resize_width else emu_to_inches(original_width)
    effective_height = emu_to_inches(original_height)

    # Check for overflow if no scale factor provided
    was_scaled = False
    if font_scale == 1.0:
        overflow, scale = estimate_content_overflow(content, effective_width, effective_height)
        if overflow:
            font_scale = scale
            was_scaled = True

    # Resize width if needed - must set ALL FOUR properties together
    if resize_width:
        body_shape.left = original_left
        body_shape.top = original_top
        body_shape.width = Inches(resize_width)
        body_shape.height = original_height

    tf = body_shape.text_frame
    tf.clear()

    # Restore all properties after clear (clear can reset auto-fit)
    body_shape.left = original_left
    body_shape.top = original_top
    body_shape.width = Inches(resize_width) if resize_width else original_width
    body_shape.height = original_height

    # Use content_items if available, otherwise fall back to content list
    items_to_render = content_items if content_items else [ContentItem(text=c) for c in content]

    for i, item in enumerate(items_to_render):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Get text and code flag from ContentItem
        item_text = item.text if isinstance(item, ContentItem) else item
        is_code = item.is_code if isinstance(item, ContentItem) else False

        # Skip figure references
        if item_text.startswith('[Figure:') or item_text.startswith('**Figure**'):
            continue

        # For code blocks, disable bullet and use monospace
        if is_code:
            level = 0
            disable_bullet(p)
            apply_formatted_text(p, item_text, is_code=True)
        else:
            # Detect indentation level and clean the text
            level, cleaned_text, is_numbered = detect_indent_level(item_text)
            p.level = level

            # For numbered items, disable the bullet
            if is_numbered:
                disable_bullet(p)

            # Apply formatted text (handles **bold**, [link](url), `code` patterns)
            apply_formatted_text(p, cleaned_text)

        # Apply scaled font size if needed
        if font_scale < 1.0:
            base_size = TEMPLATE_FONT_SIZES.get(level, 18)
            new_size = max(MIN_FONT_SIZE, int(base_size * font_scale))
            for run in p.runs:
                run.font.size = Pt(new_size)

    return was_scaled, font_scale


def insert_image(slide, image_path: Path, left: float, top: float,
                 max_width: float = None, max_height: float = None):
    """Insert an image as a shape with proper sizing."""

    if not image_path.exists():
        print(f"  Warning: Image not found: {image_path}")
        return None

    # Get image dimensions
    img_width_px, img_height_px = get_image_dimensions(image_path)
    aspect_ratio = img_width_px / img_height_px

    # Calculate size maintaining aspect ratio
    if max_width and max_height:
        # Fit within bounds
        if aspect_ratio > max_width / max_height:
            width = max_width
            height = width / aspect_ratio
        else:
            height = max_height
            width = height * aspect_ratio
    elif max_width:
        width = max_width
        height = width / aspect_ratio
    elif max_height:
        height = max_height
        width = height * aspect_ratio
    else:
        # Default size
        width = 5.0
        height = width / aspect_ratio

    # Add the picture
    pic = slide.shapes.add_picture(
        str(image_path),
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height)
    )

    return pic


def get_layout_for_slide_type(
    prs,
    slide_type: str,
    layout_mapping: LayoutMapping
) -> Tuple[object, Optional[LayoutCapabilities], Dict]:
    """
    Get the appropriate layout for a slide type using dynamic discovery.

    Args:
        prs: Presentation object
        slide_type: Type of slide (title, text_only, text_image, etc.)
        layout_mapping: LayoutMapping from discover_layouts()

    Returns:
        Tuple of (layout, LayoutCapabilities or None, fallback_config dict)
    """
    # Use dynamic layout discovery
    caps, fallback_config = layout_mapping.get_layout_for_type(slide_type)

    if caps is None:
        # No suitable layout found - use first available layout
        print(f"  Warning: No layout found for slide type '{slide_type}', using first layout")
        layout = prs.slide_layouts[0]
        return layout, None, fallback_config

    # Return the layout from the presentation by index
    layout = prs.slide_layouts[caps.index]
    return layout, caps, fallback_config


def generate_presentation(
    template_path: Path,
    slides_data: List[SlideContent],
    figures_dir: Path,
    output_path: Path,
    layout_config: LayoutConfig
):
    """Generate the presentation with text and images."""

    prs = Presentation(str(template_path))

    # Get slide dimensions
    slide_width = emu_to_inches(prs.slide_width)
    slide_height = emu_to_inches(prs.slide_height)

    print(f"Slide dimensions: {slide_width:.2f}\" x {slide_height:.2f}\"")
    print(f"Available layouts: {len(prs.slide_layouts)}")
    for i, layout in enumerate(prs.slide_layouts):
        print(f"  {i}: {layout.name}")
    print()

    # Discover layouts dynamically
    layout_mapping = discover_layouts(prs, verbose=True)

    print(f"Processing {len(slides_data)} slides...")
    print()

    # Clear existing slides if template has content slides
    # We'll create fresh slides from layouts
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    print(f"Creating {len(slides_data)} new slides from layouts...")

    image_count = 0
    diagram_count = 0
    table_count = 0

    # Track slides where font scaling was applied
    compressions = []  # List of dicts with slide info and scale factor

    for i, slide_content in enumerate(slides_data):
        # Determine effective slide type (may adjust based on content)
        effective_slide_type = slide_content.slide_type

        # For tables, determine full-width vs split layout before creating slide
        if slide_content.slide_type == 'text_table' and slide_content.table_spec:
            has_content = bool([c for c in slide_content.content if c.strip()])
            num_cols = len(slide_content.table_spec.get('headers', []))
            use_full_width = not has_content or num_cols > 3
            if use_full_width:
                effective_slide_type = 'text_table_full'
            else:
                effective_slide_type = 'text_table'

        # Create a new slide with appropriate layout using dynamic discovery
        layout, caps, fallback_config = get_layout_for_slide_type(
            prs, effective_slide_type, layout_mapping
        )
        slide = prs.slides.add_slide(layout)

        # Fill title
        if slide_content.title:
            fill_title(slide, slide_content.title)

        # Handle based on slide type
        if slide_content.slide_type == 'text_image' and slide_content.figure_path:
            # Side-by-side layout: text on left, image on right
            clean_content = [
                c for c in slide_content.content
                if not c.startswith('[Figure:') and
                   not c.startswith('**Figure**') and
                   'status-tab' not in c.lower() and
                   '.png]' not in c
            ]
            # Filter content_items similarly
            clean_content_items = [
                item for item in slide_content.content_items
                if not item.text.startswith('[Figure:') and
                   not item.text.startswith('**Figure**') and
                   'status-tab' not in item.text.lower() and
                   '.png]' not in item.text
            ] if slide_content.content_items else None

            # Determine placeholder indices from capabilities
            if caps and caps.has_typed_placeholders and caps.body_idx is not None:
                # Use typed placeholder (BODY for text)
                text_placeholder_idx = caps.body_idx
            elif caps and caps.left_content_idx is not None:
                # Use left content placeholder
                text_placeholder_idx = caps.left_content_idx
            else:
                # Fallback
                text_placeholder_idx = 1

            if caps and caps.right_content_idx is not None:
                image_placeholder_idx = caps.right_content_idx
            else:
                image_placeholder_idx = 2

            # Fill text placeholder
            was_scaled, scale = fill_body_placeholder(
                slide, clean_content, placeholder_idx=text_placeholder_idx,
                content_items=clean_content_items
            )
            if was_scaled:
                compressions.append({
                    'slide': i + 1,
                    'title': slide_content.title[:40] if slide_content.title else 'Untitled',
                    'scale': scale,
                    'base_font': 24,
                    'scaled_font': int(24 * scale)
                })

            # Resolve figure path
            fig_path = slide_content.figure_path
            if fig_path.startswith('../'):
                fig_path = fig_path[3:]  # Remove ../

            image_path = figures_dir.parent / fig_path
            if not image_path.exists():
                image_path = figures_dir / Path(fig_path).name

            if image_path.exists():
                # Insert image into OBJECT placeholder
                pic = insert_image_into_placeholder(slide, image_path, placeholder_idx=image_placeholder_idx)

                if pic:
                    image_count += 1
                    actual_left = emu_to_inches(pic.left)
                    actual_top = emu_to_inches(pic.top)
                    actual_width = emu_to_inches(pic.width)
                    actual_height = emu_to_inches(pic.height)
                    print(f"Slide {i+1}: {slide_content.title[:40]}...")
                    print(f"  Image: {image_path.name}")
                    print(f"  Position: ({actual_left:.2f}\", {actual_top:.2f}\") "
                          f"Size: {actual_width:.2f}\" x {actual_height:.2f}\"")
            else:
                print(f"Slide {i+1}: Image not found: {fig_path}")

        elif slide_content.slide_type == 'text_image_top' and slide_content.figure_path:
            # Stacked layout: image on top, text on bottom
            clean_content = [
                c for c in slide_content.content
                if not c.startswith('[Figure:') and
                   not c.startswith('**Figure**') and
                   'status-tab' not in c.lower() and
                   '.png]' not in c
            ]
            # Filter content_items similarly
            clean_content_items = [
                item for item in slide_content.content_items
                if not item.text.startswith('[Figure:') and
                   not item.text.startswith('**Figure**') and
                   'status-tab' not in item.text.lower() and
                   '.png]' not in item.text
            ] if slide_content.content_items else None

            # Determine placeholder indices from capabilities
            if caps and caps.top_content_idx is not None:
                image_placeholder_idx = caps.top_content_idx
            else:
                image_placeholder_idx = 1

            if caps and caps.has_typed_placeholders and caps.body_idx is not None:
                # Use typed placeholder (BODY for text)
                text_placeholder_idx = caps.body_idx
            elif caps and caps.bottom_content_idx is not None:
                text_placeholder_idx = caps.bottom_content_idx
            else:
                text_placeholder_idx = 2

            # Resolve figure path
            fig_path = slide_content.figure_path
            if fig_path.startswith('../'):
                fig_path = fig_path[3:]  # Remove ../

            image_path = figures_dir.parent / fig_path
            if not image_path.exists():
                image_path = figures_dir / Path(fig_path).name

            if image_path.exists():
                # Insert image into TOP placeholder
                pic = insert_image_into_placeholder(slide, image_path, placeholder_idx=image_placeholder_idx)

                if pic:
                    image_count += 1
                    actual_left = emu_to_inches(pic.left)
                    actual_top = emu_to_inches(pic.top)
                    actual_width = emu_to_inches(pic.width)
                    actual_height = emu_to_inches(pic.height)
                    print(f"Slide {i+1}: {slide_content.title[:40]}...")
                    print(f"  Image (top): {image_path.name}")
                    print(f"  Position: ({actual_left:.2f}\", {actual_top:.2f}\") "
                          f"Size: {actual_width:.2f}\" x {actual_height:.2f}\"")
            else:
                print(f"Slide {i+1}: Image not found: {fig_path}")

            # Fill BOTTOM placeholder with text
            was_scaled, scale = fill_body_placeholder(
                slide, clean_content, placeholder_idx=text_placeholder_idx,
                content_items=clean_content_items
            )
            if was_scaled:
                compressions.append({
                    'slide': i + 1,
                    'title': slide_content.title[:40] if slide_content.title else 'Untitled',
                    'scale': scale,
                    'base_font': 24,
                    'scaled_font': int(24 * scale)
                })

        elif slide_content.slide_type == 'text_diagram' and slide_content.diagram_spec:
            # Fill content with resized width (same as images)
            was_scaled, scale = fill_body(
                slide, slide_content.content, resize_width=layout_config.text_width,
                content_items=slide_content.content_items
            )
            if was_scaled:
                compressions.append({
                    'slide': i + 1,
                    'title': slide_content.title[:40] if slide_content.title else 'Untitled',
                    'scale': scale,
                    'base_font': 24,
                    'scaled_font': int(24 * scale)
                })

            # Configure diagram position
            diagram_config = DiagramConfig(
                left=layout_config.image_left,
                top=layout_config.image_top,
                width=layout_config.image_width,
                height=slide_height - layout_config.image_top - 0.5
            )

            # Render the diagram
            shapes = render_diagram(slide, slide_content.diagram_spec, diagram_config)

            if shapes:
                diagram_count += 1
                diagram_type = slide_content.diagram_spec.get('type', 'unknown')
                node_count = len(slide_content.diagram_spec.get('nodes', []))
                print(f"Slide {i+1}: {slide_content.title[:40]}...")
                print(f"  Diagram: {diagram_type} with {node_count} nodes, {len(shapes)} shapes")

        elif slide_content.slide_type == 'text_table' and slide_content.table_spec:
            # Determine if table should be full-width or split layout
            # Full-width when: no bullet content, or table has many columns (>3)
            has_content = bool([c for c in slide_content.content if c.strip()])
            num_cols = len(slide_content.table_spec.get('headers', []))
            use_full_width = not has_content or num_cols > 3

            table_shape = None

            if use_full_width:
                # Full-width table - use body placeholder for table
                if caps and caps.body_idx is not None:
                    table_shape = insert_table_into_placeholder(
                        slide, slide_content.table_spec, placeholder_idx=caps.body_idx
                    )
                else:
                    # Fallback to manual positioning
                    table_config = TableConfig(left=0.5, top=1.5, width=12.3)
                    hide_content_placeholders(slide)
                    table_shape = render_table(slide, slide_content.table_spec, table_config)
            else:
                # Split layout - text on left, table on right
                # Determine placeholder indices from capabilities
                if caps and caps.has_typed_placeholders and caps.body_idx is not None:
                    text_placeholder_idx = caps.body_idx
                elif caps and caps.left_content_idx is not None:
                    text_placeholder_idx = caps.left_content_idx
                else:
                    text_placeholder_idx = 1

                if caps and caps.right_content_idx is not None:
                    table_placeholder_idx = caps.right_content_idx
                else:
                    table_placeholder_idx = 2

                # Fill text placeholder
                was_scaled, scale = fill_body_placeholder(
                    slide, slide_content.content, placeholder_idx=text_placeholder_idx,
                    content_items=slide_content.content_items
                )
                if was_scaled:
                    compressions.append({
                        'slide': i + 1,
                        'title': slide_content.title[:40] if slide_content.title else 'Untitled',
                        'scale': scale,
                        'base_font': 24,
                        'scaled_font': int(24 * scale)
                    })
                # Insert table into right placeholder
                table_shape = insert_table_into_placeholder(
                    slide, slide_content.table_spec, placeholder_idx=table_placeholder_idx
                )

            if table_shape:
                table_count += 1
                headers = slide_content.table_spec.get('headers', [])
                rows = slide_content.table_spec.get('rows', [])
                print(f"Slide {i+1}: {slide_content.title[:40]}...")
                print(f"  Table: {len(headers)} cols x {len(rows)+1} rows")

        elif slide_content.slide_type in ('text_only', 'section'):
            # Full-width text
            was_scaled, scale = fill_body(
                slide, slide_content.content,
                content_items=slide_content.content_items
            )
            if was_scaled:
                compressions.append({
                    'slide': i + 1,
                    'title': slide_content.title[:40] if slide_content.title else 'Untitled',
                    'scale': scale,
                    'base_font': 24,
                    'scaled_font': int(24 * scale)
                })

        elif slide_content.slide_type == 'title':
            # Title slide - use content as subtitle
            if slide_content.content:
                # Join content items as subtitle (usually includes date)
                subtitle_text = '\n'.join(slide_content.content)
                fill_subtitle(slide, subtitle_text)

    # Save
    prs.save(str(output_path))

    # Fix view mode to open in Normal view (not Slide Master view)
    fix_view_mode(output_path)

    print()
    print(f"Generated {len(slides_data)} slides with {image_count} images, {diagram_count} diagrams, and {table_count} tables")
    print(f"Saved to: {output_path}")

    # Print compression summary if any slides were scaled
    if compressions:
        print()
        print("=" * 60)
        print("CONTENT COMPRESSION SUMMARY")
        print("=" * 60)
        for c in compressions:
            print(f"  Slide {c['slide']}: {c['title']}...")
            print(f"    Font scaled to {c['scale']*100:.0f}% ({c['base_font']}pt -> {c['scaled_font']}pt)")
        print()
        print("Consider editing content on these slides if text is too small.")


def fix_view_mode(pptx_path: Path):
    """
    Fix the presentation to open in Normal view instead of Slide Master view.

    This modifies the viewProps.xml inside the pptx to set lastView="sldView".
    """
    import zipfile
    import tempfile
    import shutil

    pptx_path = Path(pptx_path)

    # Create a temporary file
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
        tmp_path = tmp.name

    try:
        with zipfile.ZipFile(pptx_path, 'r') as zip_in:
            with zipfile.ZipFile(tmp_path, 'w', zipfile.ZIP_DEFLATED) as zip_out:
                for item in zip_in.namelist():
                    data = zip_in.read(item)

                    # Modify viewProps.xml to set Normal view
                    if item == 'ppt/viewProps.xml':
                        content = data.decode('utf-8')
                        # Change lastView to sldView (Normal/Slide view)
                        if 'lastView="sldMasterView"' in content:
                            content = content.replace('lastView="sldMasterView"', 'lastView="sldView"')
                        elif 'lastView=' not in content:
                            # Add lastView attribute if not present
                            content = content.replace('<p:viewPr', '<p:viewPr lastView="sldView"')
                        data = content.encode('utf-8')

                    zip_out.writestr(item, data)

        # Replace original with modified
        shutil.move(tmp_path, pptx_path)
    except Exception as e:
        # Clean up temp file on error
        if Path(tmp_path).exists():
            Path(tmp_path).unlink()
        print(f"Warning: Could not fix view mode: {e}")


def main():
    parser = argparse.ArgumentParser(
        description="Generate PowerPoint from PRESENTATION.md with images",
        formatter_class=argparse.RawDescriptionHelpFormatter
    )
    parser.add_argument(
        "--template", "-t",
        default=None,
        help=f"PowerPoint template file (default: {DEFAULT_TEMPLATE.name})"
    )
    parser.add_argument(
        "--source", "-s",
        required=True,
        help="PRESENTATION.md source file"
    )
    parser.add_argument(
        "--figures", "-f",
        default=None,
        help="Figures directory (default: figures/ relative to source)"
    )
    parser.add_argument(
        "--output", "-o",
        required=True,
        help="Output PowerPoint file"
    )
    parser.add_argument(
        "--text-width",
        type=float,
        default=6.5,
        help="Width of text area for image slides (default: 6.5)"
    )
    parser.add_argument(
        "--image-left",
        type=float,
        default=7.0,
        help="Left position of images (default: 7.0)"
    )
    parser.add_argument(
        "--image-width",
        type=float,
        default=5.5,
        help="Max width of images (default: 5.5)"
    )
    parser.add_argument(
        "--image-top",
        type=float,
        default=1.5,
        help="Top position of images (default: 1.5)"
    )

    args = parser.parse_args()

    # Use default template if not specified
    if args.template:
        template_path = Path(args.template)
    else:
        template_path = DEFAULT_TEMPLATE
        print(f"Using default template: {template_path}")

    source_path = Path(args.source)
    output_path = Path(args.output)

    # Default figures directory to figures/ relative to source
    if args.figures:
        figures_dir = Path(args.figures)
    else:
        figures_dir = source_path.parent / "figures"
        if not figures_dir.exists():
            # Fall back to source directory itself
            figures_dir = source_path.parent

    if not template_path.exists():
        print(f"Error: Template not found: {template_path}")
        if template_path == DEFAULT_TEMPLATE:
            print("  The default template is missing. Please ensure templates/cellularintelligence.pptx exists.")
        sys.exit(1)

    if not source_path.exists():
        print(f"Error: Source not found: {source_path}")
        sys.exit(1)

    # Parse PRESENTATION.md
    print(f"Parsing: {source_path}")
    slides_data = parse_presentation_md(source_path)
    print(f"Found {len(slides_data)} slides")

    # Count slides with figures
    image_slides = [s for s in slides_data if s.figure_path]
    print(f"Slides with figures: {len(image_slides)}")
    for s in image_slides:
        print(f"  Slide {s.number}: {s.figure_path}")

    # Count slides with diagrams
    diagram_slides = [s for s in slides_data if s.diagram_spec]
    print(f"Slides with diagrams: {len(diagram_slides)}")
    for s in diagram_slides:
        dtype = s.diagram_spec.get('type', 'unknown')
        nodes = len(s.diagram_spec.get('nodes', []))
        print(f"  Slide {s.number}: {dtype} diagram ({nodes} nodes)")

    # Count slides with tables
    table_slides = [s for s in slides_data if s.table_spec]
    print(f"Slides with tables: {len(table_slides)}")
    for s in table_slides:
        headers = s.table_spec.get('headers', [])
        rows = s.table_spec.get('rows', [])
        print(f"  Slide {s.number}: {len(headers)} cols x {len(rows)+1} rows")
    print()

    # Configure layouts
    layout_config = LayoutConfig(
        text_width=args.text_width,
        image_left=args.image_left,
        image_width=args.image_width,
        image_top=args.image_top
    )

    # Generate
    generate_presentation(
        template_path,
        slides_data,
        figures_dir,
        output_path,
        layout_config
    )


if __name__ == "__main__":
    main()
