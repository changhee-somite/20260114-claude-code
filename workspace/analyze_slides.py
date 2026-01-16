#!/usr/bin/env python3
"""
Analyze slides in a PowerPoint presentation.

This script provides information about slides, their layouts, and existing shapes
to help plan image insertion positions.

Usage:
    python analyze_slides.py <pptx_file> [--slide N]

Examples:
    python analyze_slides.py presentation.pptx
    python analyze_slides.py presentation.pptx --slide 5
"""

import argparse
import sys
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE


def emu_to_inches(emu):
    """Convert EMUs to inches."""
    return emu / 914400


def analyze_presentation(pptx_path, specific_slide=None):
    """Analyze a presentation and print slide information."""
    prs = Presentation(pptx_path)

    # Presentation dimensions
    slide_width = emu_to_inches(prs.slide_width)
    slide_height = emu_to_inches(prs.slide_height)

    print(f"=" * 60)
    print(f"Presentation: {pptx_path}")
    print(f"=" * 60)
    print(f"Dimensions: {slide_width:.2f}\" x {slide_height:.2f}\" (W x H)")
    print(f"Total slides: {len(prs.slides)}")
    print()

    # Available layouts
    print("Available Slide Layouts:")
    print("-" * 40)
    for i, layout in enumerate(prs.slide_layouts):
        print(f"  Layout {i}: {layout.name}")
    print()

    # Analyze slides
    slides_to_analyze = [specific_slide - 1] if specific_slide else range(len(prs.slides))

    for slide_idx in slides_to_analyze:
        if slide_idx < 0 or slide_idx >= len(prs.slides):
            print(f"Slide {slide_idx + 1} does not exist")
            continue

        slide = prs.slides[slide_idx]
        layout_name = slide.slide_layout.name if slide.slide_layout else "Unknown"

        print(f"Slide {slide_idx + 1}: Layout = '{layout_name}'")
        print("-" * 40)

        # Find content bounds to suggest image placement
        content_shapes = []

        for shape in slide.shapes:
            left = emu_to_inches(shape.left)
            top = emu_to_inches(shape.top)
            width = emu_to_inches(shape.width)
            height = emu_to_inches(shape.height)

            shape_type = shape.shape_type

            # Get text preview if available
            text_preview = ""
            if shape.has_text_frame:
                full_text = shape.text_frame.text[:50]
                if full_text:
                    text_preview = f' "{full_text}..."' if len(shape.text_frame.text) > 50 else f' "{full_text}"'

            print(f"  {shape.name}")
            print(f"    Type: {shape_type}")
            print(f"    Position: ({left:.2f}\", {top:.2f}\")")
            print(f"    Size: {width:.2f}\" x {height:.2f}\"{text_preview}")

            content_shapes.append({
                'name': shape.name,
                'left': left,
                'top': top,
                'width': width,
                'height': height,
                'bottom': top + height,
                'right': left + width
            })

        # Suggest image placement areas
        if content_shapes:
            max_bottom = max(s['bottom'] for s in content_shapes)
            max_right = max(s['right'] for s in content_shapes)

            print()
            print("  Suggested image areas:")

            # Below content
            if max_bottom < slide_height - 1:
                available_height = slide_height - max_bottom - 0.5
                print(f"    Below content: top={max_bottom + 0.25:.2f}\", height={available_height:.2f}\" available")

            # Right side (for text-left/image-right layout)
            if max_right < slide_width * 0.6:
                available_width = slide_width - max_right - 0.5
                print(f"    Right side: left={max_right + 0.25:.2f}\", width={available_width:.2f}\" available")

        print()


def main():
    parser = argparse.ArgumentParser(
        description="Analyze PowerPoint slides for image insertion planning"
    )
    parser.add_argument("pptx_file", help="Path to the PowerPoint file")
    parser.add_argument("--slide", "-s", type=int, help="Analyze specific slide number (1-indexed)")

    args = parser.parse_args()

    try:
        analyze_presentation(args.pptx_file, args.slide)
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
