#!/usr/bin/env python3
"""
Insert images into PowerPoint slides with CONTENT-AWARE positioning.

This improved version analyzes existing slide content to avoid overlaps.

Key improvements over v1:
- Analyzes existing text shapes before placing images
- Smart presets that adapt to actual content bounds
- Overlap detection and warnings
- Option to resize text placeholders to make room

Usage:
    python insert_images_v2.py <input_pptx> <image_mapping_json> <output_pptx>

Example:
    python insert_images_v2.py working.pptx image-mapping.json final.pptx
"""

import argparse
import json
import os
import sys
from pathlib import Path
from dataclasses import dataclass
from typing import List, Optional, Tuple

from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image


@dataclass
class BoundingBox:
    """Represents a rectangular area in inches."""
    left: float
    top: float
    width: float
    height: float

    @property
    def right(self) -> float:
        return self.left + self.width

    @property
    def bottom(self) -> float:
        return self.top + self.height

    def overlaps(self, other: 'BoundingBox') -> bool:
        """Check if this box overlaps with another."""
        h_overlap = not (self.right <= other.left or self.left >= other.right)
        v_overlap = not (self.bottom <= other.top or self.top >= other.bottom)
        return h_overlap and v_overlap


def emu_to_inches(emu) -> float:
    """Convert EMUs to inches."""
    return emu / 914400


def inches_to_emu(inches) -> int:
    """Convert inches to EMUs."""
    return int(inches * 914400)


def get_image_dimensions(image_path) -> Tuple[int, int]:
    """Get image dimensions in pixels."""
    with Image.open(image_path) as img:
        return img.size  # (width, height)


def analyze_slide_content(slide, slide_width: float, slide_height: float) -> dict:
    """
    Analyze slide content to find text bounds and available space.

    Returns dict with:
        - text_bounds: BoundingBox of all text content
        - title_bounds: BoundingBox of title (if exists)
        - content_bounds: BoundingBox of main content area
        - available_right: Space available on right side
        - available_bottom: Space available at bottom
    """
    text_shapes = []
    title_shape = None

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            continue

        if not hasattr(shape, 'text_frame') or not shape.has_text_frame:
            continue

        text = shape.text_frame.text.strip()
        if not text:
            continue

        box = BoundingBox(
            left=emu_to_inches(shape.left),
            top=emu_to_inches(shape.top),
            width=emu_to_inches(shape.width),
            height=emu_to_inches(shape.height)
        )

        # Check if this is a title (usually at top, spans width)
        if box.top < 1.5 and 'title' in shape.name.lower():
            title_shape = box
        else:
            text_shapes.append(box)

    # Calculate combined text bounds
    if text_shapes:
        all_shapes = text_shapes + ([title_shape] if title_shape else [])
        text_bounds = BoundingBox(
            left=min(s.left for s in all_shapes),
            top=min(s.top for s in all_shapes),
            width=max(s.right for s in all_shapes) - min(s.left for s in all_shapes),
            height=max(s.bottom for s in all_shapes) - min(s.top for s in all_shapes)
        )

        # Content bounds (excluding title)
        content_bounds = BoundingBox(
            left=min(s.left for s in text_shapes),
            top=min(s.top for s in text_shapes),
            width=max(s.right for s in text_shapes) - min(s.left for s in text_shapes),
            height=max(s.bottom for s in text_shapes) - min(s.top for s in text_shapes)
        )
    else:
        text_bounds = None
        content_bounds = None

    # Calculate available space
    margin = 0.25

    if content_bounds:
        available_right = slide_width - content_bounds.right - margin
        available_bottom = slide_height - content_bounds.bottom - margin
    else:
        available_right = slide_width - margin * 2
        available_bottom = slide_height - margin * 2

    return {
        'text_bounds': text_bounds,
        'title_bounds': title_shape,
        'content_bounds': content_bounds,
        'available_right': max(0, available_right),
        'available_bottom': max(0, available_bottom),
        'slide_width': slide_width,
        'slide_height': slide_height
    }


def calculate_smart_position(
    preset: str,
    analysis: dict,
    img_width_px: int,
    img_height_px: int
) -> Tuple[float, float, float, float]:
    """
    Calculate image position based on preset and ACTUAL slide content.

    Returns (left, top, width, height) in inches.
    """
    slide_width = analysis['slide_width']
    slide_height = analysis['slide_height']
    content_bounds = analysis['content_bounds']
    title_bounds = analysis['title_bounds']

    margin = 0.25
    img_ratio = img_width_px / img_height_px

    if preset == "right_of_content":
        # Place image to the right of actual content
        if content_bounds and content_bounds.right < slide_width * 0.7:
            # Content leaves room on the right
            available_left = content_bounds.right + margin
            available_width = slide_width - available_left - margin

            # Vertical: below title if exists, otherwise from top
            if title_bounds:
                available_top = title_bounds.bottom + margin
            else:
                available_top = margin
            available_height = slide_height - available_top - margin
        else:
            # Content is full width, fall back to right half overlay
            available_left = slide_width * 0.55
            available_width = slide_width * 0.4
            available_top = 1.5 if title_bounds else margin
            available_height = slide_height - available_top - margin

        # Calculate size maintaining aspect ratio
        if img_ratio > available_width / available_height:
            width = available_width
            height = width / img_ratio
        else:
            height = available_height
            width = height * img_ratio

        # Center in available space
        left = available_left + (available_width - width) / 2
        top = available_top + (available_height - height) / 2

        return left, top, width, height

    elif preset == "right_of_content_top":
        # Place image to the right of content, TOP-ALIGNED (not centered)
        if content_bounds and content_bounds.right < slide_width * 0.7:
            available_left = content_bounds.right + margin
            available_width = slide_width - available_left - margin

            # Top aligned with content (or below title)
            if title_bounds:
                available_top = title_bounds.bottom + margin
            elif content_bounds:
                available_top = content_bounds.top
            else:
                available_top = margin
            available_height = slide_height - available_top - margin
        else:
            available_left = slide_width * 0.55
            available_width = slide_width * 0.4
            available_top = 1.5 if title_bounds else margin
            available_height = slide_height - available_top - margin

        # Calculate size maintaining aspect ratio
        if img_ratio > available_width / available_height:
            width = available_width
            height = width / img_ratio
        else:
            height = available_height
            width = height * img_ratio

        # TOP-ALIGNED: center horizontally, align to top
        left = available_left + (available_width - width) / 2
        top = available_top  # <-- Top aligned, not centered

        return left, top, width, height

    elif preset == "below_content":
        # Place image below actual content
        if content_bounds:
            available_top = content_bounds.bottom + margin
        elif title_bounds:
            available_top = title_bounds.bottom + margin
        else:
            available_top = slide_height * 0.5

        available_height = slide_height - available_top - margin
        available_width = slide_width - margin * 2

        if img_ratio > available_width / available_height:
            width = available_width
            height = width / img_ratio
        else:
            height = available_height
            width = height * img_ratio

        left = (slide_width - width) / 2
        top = available_top

        return left, top, width, height

    elif preset == "right_half_safe":
        # Right half, but respects title area
        if title_bounds:
            available_top = title_bounds.bottom + margin
        else:
            available_top = margin

        available_left = slide_width / 2 + margin
        available_width = slide_width / 2 - margin * 2
        available_height = slide_height - available_top - margin

        if img_ratio > available_width / available_height:
            width = available_width
            height = width / img_ratio
        else:
            height = available_height
            width = height * img_ratio

        left = available_left + (available_width - width) / 2
        top = available_top + (available_height - height) / 2

        return left, top, width, height

    elif preset == "center_below_title":
        # Centered but below title
        if title_bounds:
            available_top = title_bounds.bottom + margin
        else:
            available_top = 1.2  # Assume title area

        available_height = slide_height - available_top - margin
        available_width = slide_width * 0.8

        if img_ratio > available_width / available_height:
            width = available_width
            height = width / img_ratio
        else:
            height = available_height
            width = height * img_ratio

        left = (slide_width - width) / 2
        top = available_top + (available_height - height) / 2

        return left, top, width, height

    elif preset == "background":
        # Full slide background (behind everything)
        return 0, 0, slide_width, slide_height

    else:
        raise ValueError(f"Unknown preset: {preset}. Use: right_of_content, below_content, right_half_safe, center_below_title, background")


def insert_images(input_pptx, mapping_json, output_pptx, base_path=None, resize_text=False):
    """
    Insert images with content-aware positioning.
    """
    prs = Presentation(input_pptx)
    slide_width = emu_to_inches(prs.slide_width)
    slide_height = emu_to_inches(prs.slide_height)

    with open(mapping_json, 'r') as f:
        mapping = json.load(f)

    if base_path is None:
        base_path = Path(mapping_json).parent
    else:
        base_path = Path(base_path)

    inserted_count = 0
    warnings = []
    errors = []

    for slide_config in mapping.get('slides', []):
        slide_number = slide_config.get('slide_number')

        if slide_number is None:
            errors.append("Missing 'slide_number' in slide configuration")
            continue

        slide_idx = slide_number - 1

        if slide_idx < 0 or slide_idx >= len(prs.slides):
            errors.append(f"Slide {slide_number} does not exist")
            continue

        slide = prs.slides[slide_idx]

        # Analyze slide content BEFORE inserting images
        analysis = analyze_slide_content(slide, slide_width, slide_height)

        for img_config in slide_config.get('images', []):
            image_path = img_config.get('path')

            if not image_path:
                errors.append(f"Slide {slide_number}: Missing 'path'")
                continue

            full_path = base_path / image_path if not os.path.isabs(image_path) else Path(image_path)

            if not full_path.exists():
                errors.append(f"Slide {slide_number}: Image not found: {full_path}")
                continue

            try:
                img_width_px, img_height_px = get_image_dimensions(full_path)

                if 'preset' in img_config:
                    left, top, width, height = calculate_smart_position(
                        img_config['preset'],
                        analysis,
                        img_width_px,
                        img_height_px
                    )
                else:
                    # Manual positioning
                    left = img_config.get('left', 0)
                    top = img_config.get('top', 0)
                    width = img_config.get('width')
                    height = img_config.get('height')

                    if width is None and height is None:
                        width = img_width_px / 96
                        height = img_height_px / 96
                    elif width is None:
                        width = height * (img_width_px / img_height_px)
                    elif height is None:
                        height = width * (img_height_px / img_width_px)

                # Check for overlaps with text
                img_box = BoundingBox(left, top, width, height)

                if analysis['content_bounds']:
                    if img_box.overlaps(analysis['content_bounds']):
                        warnings.append(
                            f"Slide {slide_number}: Image may overlap with text content. "
                            f"Image: L={left:.1f}\" R={left+width:.1f}\" | "
                            f"Text: L={analysis['content_bounds'].left:.1f}\" R={analysis['content_bounds'].right:.1f}\""
                        )

                # Optionally resize text placeholder to make room
                if resize_text and analysis['content_bounds']:
                    for shape in slide.shapes:
                        if (hasattr(shape, 'text_frame') and
                            shape.has_text_frame and
                            'placeholder' in shape.name.lower() and
                            'title' not in shape.name.lower()):

                            current_right = emu_to_inches(shape.left + shape.width)
                            if current_right > left - 0.25:
                                new_width = left - emu_to_inches(shape.left) - 0.25
                                if new_width > 2:  # Minimum width
                                    shape.width = inches_to_emu(new_width)
                                    print(f"  Resized text placeholder to {new_width:.1f}\" wide")

                # Add the image
                slide.shapes.add_picture(
                    str(full_path),
                    Inches(left),
                    Inches(top),
                    Inches(width),
                    Inches(height)
                )

                inserted_count += 1
                print(f"  Slide {slide_number}: {image_path}")
                print(f"    Position: ({left:.2f}\", {top:.2f}\") Size: {width:.2f}\" x {height:.2f}\"")
                if analysis['content_bounds']:
                    print(f"    Content ends at: {analysis['content_bounds'].right:.2f}\"")

            except Exception as e:
                errors.append(f"Slide {slide_number}: Error - {e}")

    # Report
    print()
    print(f"Inserted {inserted_count} image(s)")

    if warnings:
        print(f"\n⚠️  Warnings ({len(warnings)}):")
        for w in warnings:
            print(f"  - {w}")

    if errors:
        print(f"\n❌ Errors ({len(errors)}):")
        for e in errors:
            print(f"  - {e}")

    prs.save(output_pptx)
    print(f"\n✅ Saved to: {output_pptx}")

    return len(errors) == 0


def main():
    parser = argparse.ArgumentParser(
        description="Insert images with content-aware positioning",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Smart Presets (v2):
  - right_of_content: Places image to the RIGHT of actual text content
  - below_content: Places image BELOW actual text content
  - right_half_safe: Right half, but avoids title area
  - center_below_title: Centered, but below title
  - background: Full slide background

Example mapping:
{
  "slides": [
    {
      "slide_number": 5,
      "images": [{"path": "fig.png", "preset": "right_of_content"}]
    }
  ]
}
        """
    )
    parser.add_argument("input_pptx", help="Input PowerPoint file")
    parser.add_argument("mapping_json", help="Image mapping JSON")
    parser.add_argument("output_pptx", help="Output PowerPoint file")
    parser.add_argument("--resize-text", "-r", action="store_true",
                       help="Resize text placeholders to make room for images")
    parser.add_argument("--base-path", "-b", help="Base path for image files")

    args = parser.parse_args()

    print(f"Input: {args.input_pptx}")
    print(f"Mapping: {args.mapping_json}")
    print(f"Output: {args.output_pptx}")
    print()

    success = insert_images(
        args.input_pptx,
        args.mapping_json,
        args.output_pptx,
        args.base_path,
        args.resize_text
    )

    sys.exit(0 if success else 1)


if __name__ == "__main__":
    main()
